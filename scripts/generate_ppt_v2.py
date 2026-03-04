# -*- coding: utf-8 -*-
"""
动态数据管理平台 — 系统功能介绍 PPT V2 生成脚本
15 页完整版，覆盖全部功能模块
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 全局配色 ──────────────────────────────────────────────
PRIMARY = RGBColor(0x40, 0x9E, 0xFF)
PRIMARY_DARK = RGBColor(0x33, 0x7E, 0xCC)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x67, 0xC2, 0x3A)
ACCENT2 = RGBColor(0xE6, 0xA2, 0x3C)
ACCENT3 = RGBColor(0xF5, 0x6C, 0x6C)
PURPLE = RGBColor(0x90, 0x9B, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xC0, 0xC4, 0xCC)
VERY_LIGHT = RGBColor(0xE4, 0xE7, 0xED)
CARD_BG = RGBColor(0x23, 0x2D, 0x4A)
BORDER = RGBColor(0x30, 0x3D, 0x5F)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TOTAL = 15


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.08
    return shape


def add_flat_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txbox


def add_multiline(slide, left, top, width, height, lines, size=12, color=LIGHT_TEXT, spacing=0.32):
    """添加多行文本，每行独立"""
    y = top
    for line in lines:
        add_text(slide, left, y, width, Inches(spacing), line, size=size, color=color)
        y += Inches(spacing)


def add_top_bar(slide):
    add_flat_rect(slide, 0, 0, SLIDE_WIDTH, Inches(0.06), PRIMARY)


def add_page_num(slide, num):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num} / {TOTAL}", size=11, color=LIGHT_TEXT, align=PP_ALIGN.RIGHT)


def add_slide_header(slide, title, subtitle, page_num):
    add_top_bar(slide)
    add_page_num(slide, page_num)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             title, size=32, color=WHITE, bold=True)
    add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             subtitle, size=14, color=LIGHT_TEXT)


def add_tag(slide, x, y, text, color=PRIMARY, bg=None):
    w = Inches(len(text) * 0.13 + 0.4)
    add_rect(slide, x, y, w, Inches(0.34), bg or RGBColor(0x2A, 0x3A, 0x5C),
             border_color=color, border_width=Pt(1))
    add_text(slide, x + Inches(0.05), y + Inches(0.02), w - Inches(0.1), Inches(0.3),
             text, size=10, color=color, align=PP_ALIGN.CENTER)
    return w


# ══════════════════════════════════════════════════════════
# 第 1 页 - 封面
# ══════════════════════════════════════════════════════════
def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_circle(slide, Inches(10.0), Inches(-1.5), Inches(4.5), RGBColor(0x1E, 0x3A, 0x6E))
    add_circle(slide, Inches(11.5), Inches(0.0), Inches(3.0), RGBColor(0x24, 0x44, 0x7A))
    add_circle(slide, Inches(-1.5), Inches(5.0), Inches(3.5), RGBColor(0x1E, 0x3A, 0x6E))

    add_flat_rect(slide, 0, 0, Inches(0.12), SLIDE_HEIGHT, PRIMARY)

    add_text(slide, Inches(1.2), Inches(1.5), Inches(8), Inches(1.2),
             "动态数据管理平台", size=52, color=WHITE, bold=True)
    add_text(slide, Inches(1.2), Inches(2.8), Inches(8), Inches(0.6),
             "Dynamic Data Management Platform", size=22, color=PRIMARY)

    add_flat_rect(slide, Inches(1.2), Inches(3.6), Inches(1.5), Inches(0.04), PRIMARY)

    add_text(slide, Inches(1.2), Inches(3.9), Inches(8), Inches(0.6),
             "配置驱动  ·  灵活高效  ·  零代码扩展", size=22, color=LIGHT_TEXT)
    add_text(slide, Inches(1.2), Inches(4.5), Inches(7), Inches(0.8),
             "低代码驱动的企业级数据管理解决方案\n让非开发人员也能快速构建数据应用",
             size=16, color=LIGHT_TEXT)

    tags = ["Vue 3", "TypeScript", "Element Plus", "Flask", "PostgreSQL"]
    x = Inches(1.2)
    for tag in tags:
        w = add_tag(slide, x, Inches(5.6), tag)
        x += w + Inches(0.15)


# ══════════════════════════════════════════════════════════
# 第 2 页 - 平台概述
# ══════════════════════════════════════════════════════════
def slide_02_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "平台概述",
                     "一套配置驱动的低代码数据管理系统，无需编写代码即可创建新业务模块", 2)

    concepts = [
        ("配", "页面配置 PageConfig", "定义字段架构和验证规则\nJSONB 存储字段定义\n17 种控件类型", PRIMARY),
        ("+", "菜单入口 Menu", "将 URL 链接到 PageConfig\n动态生成侧边栏导航\n支持 1-3 级菜单树", ACCENT),
        ("=", "自动生成 UI & API", "完整 CRUD 页面\n动态表单 + 数据表格\nRESTful API 端点", ACCENT2),
    ]

    x = Inches(0.8)
    for icon, title, desc, color in concepts:
        add_rect(slide, x, Inches(1.8), Inches(3.7), Inches(2.4), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        ic = add_circle(slide, x + Inches(0.3), Inches(2.1), Inches(0.7), color)
        add_text(slide, x + Inches(0.3), Inches(2.22), Inches(0.7), Inches(0.5),
                 icon, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(1.2), Inches(2.15), Inches(2.3), Inches(0.4),
                 title, size=15, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.3), Inches(2.9), Inches(3.1), Inches(1.2),
                 desc, size=12, color=LIGHT_TEXT)
        x += Inches(3.95)

    # 对比：传统 vs 本平台
    add_text(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(0.5),
             "传统方式 vs 本平台", size=20, color=WHITE, bold=True)

    # 传统方式
    add_rect(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(2.0), CARD_BG,
             border_color=ACCENT3, border_width=Pt(1))
    add_text(slide, Inches(1.1), Inches(5.2), Inches(3), Inches(0.35),
             "❌ 传统开发方式", size=14, color=ACCENT3, bold=True)
    add_multiline(slide, Inches(1.1), Inches(5.6), Inches(5), Inches(1.5), [
        "•  每个业务需要新建数据库表 + 编写迁移脚本",
        "•  后端编写 API 路由、校验逻辑、序列化",
        "•  前端编写列表页、表单页、详情页",
        "•  周期长、成本高、修改需要重新发布",
    ], size=11)

    # 本平台
    add_rect(slide, Inches(7.0), Inches(5.1), Inches(5.5), Inches(2.0), CARD_BG,
             border_color=ACCENT, border_width=Pt(1))
    add_text(slide, Inches(7.3), Inches(5.2), Inches(3), Inches(0.35),
             "✓ 配置驱动方式", size=14, color=ACCENT, bold=True)
    add_multiline(slide, Inches(7.3), Inches(5.6), Inches(5), Inches(1.5), [
        "•  定义 PageConfig + Menu，系统自动生成一切",
        "•  无需建表、无需数据库迁移",
        "•  零代码扩展，实时生效",
        "•  修改字段配置即可调整业务逻辑",
    ], size=11, color=ACCENT)


# ══════════════════════════════════════════════════════════
# 第 3 页 - 系统架构
# ══════════════════════════════════════════════════════════
def slide_03_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "系统架构", "现代化全栈技术方案，前后端分离架构", 3)

    # 前端层
    add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.6), CARD_BG,
             border_color=PRIMARY, border_width=Pt(1.5))
    add_text(slide, Inches(1.1), Inches(1.8), Inches(3), Inches(0.4),
             "🖥  前端 Frontend", size=18, color=PRIMARY, bold=True)

    fe = [
        ("Vue 3", "Composition API + <script setup>"),
        ("TypeScript", "完整类型系统"),
        ("Pinia", "状态管理 Store"),
        ("Element Plus", "企业级 UI 组件库"),
        ("Vite", "极速 HMR 热更新"),
    ]
    y = Inches(2.3)
    for name, desc in fe:
        add_text(slide, Inches(1.3), y, Inches(1.6), Inches(0.28),
                 name, size=12, color=WHITE, bold=True)
        add_text(slide, Inches(2.9), y, Inches(3.0), Inches(0.28),
                 desc, size=11, color=LIGHT_TEXT)
        y += Inches(0.35)

    # 后端层
    add_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(2.6), CARD_BG,
             border_color=ACCENT, border_width=Pt(1.5))
    add_text(slide, Inches(7.3), Inches(1.8), Inches(3), Inches(0.4),
             "⚙  后端 Backend", size=18, color=ACCENT, bold=True)

    be = [
        ("Python Flask", "轻量 Web 框架，13+ 蓝图"),
        ("PostgreSQL", "JSONB 动态数据存储"),
        ("psycopg2", "连接池，高效 DB 交互"),
        ("JWT Auth", "Token 认证 + 角色鉴权"),
        ("Pytest", "241 后端单元测试"),
    ]
    y = Inches(2.3)
    for name, desc in be:
        add_text(slide, Inches(7.5), y, Inches(1.6), Inches(0.28),
                 name, size=12, color=WHITE, bold=True)
        add_text(slide, Inches(9.1), y, Inches(3.0), Inches(0.28),
                 desc, size=11, color=LIGHT_TEXT)
        y += Inches(0.35)

    # 数据库层
    add_text(slide, Inches(0.8), Inches(4.6), Inches(6), Inches(0.4),
             "数据库核心表", size=18, color=WHITE, bold=True)

    tables = [
        ("dynamic_data", "所有业务数据统一存储\ncollection + data JSONB", PRIMARY),
        ("page_configs", "字段架构定义\nfields JSONB 驱动 UI", ACCENT),
        ("data_relations", "M:N 关系映射\n双向关联自动同步", ACCENT2),
        ("menus", "菜单树结构\n角色级 RBAC 控制", PURPLE),
    ]

    x = Inches(0.8)
    for name, desc, color in tables:
        add_rect(slide, x, Inches(5.0), Inches(2.8), Inches(1.5), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_flat_rect(slide, x, Inches(5.0), Inches(2.8), Inches(0.05), color)
        add_text(slide, x + Inches(0.2), Inches(5.15), Inches(2.4), Inches(0.3),
                 name, size=13, color=color, bold=True, font="Consolas")
        add_text(slide, x + Inches(0.2), Inches(5.5), Inches(2.4), Inches(0.8),
                 desc, size=10, color=LIGHT_TEXT)
        x += Inches(3.05)

    add_text(slide, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
             "其他表: users · export_scripts · validation_scripts · etl_tasks · etl_logs · api_keys · operation_logs · backups · backup_settings",
             size=10, color=LIGHT_TEXT)


# ══════════════════════════════════════════════════════════
# 第 4 页 - 配置驱动
# ══════════════════════════════════════════════════════════
def slide_04_config_driven(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "核心设计：配置驱动",
                     "定义 PageConfig + Menu = 自动生成完整的数据管理页面和 API，无需编写任何代码", 4)

    # 流程步骤
    steps = [
        ("01", "定义 PageConfig", "管理员通过可视化界面\n定义字段名、类型、校验规则\n以 JSONB 存储于 page_configs", PRIMARY),
        ("02", "创建菜单入口", "将 URL 路径链接到\nPageConfig，配置菜单层级\n设置角色访问权限", ACCENT),
        ("03", "系统自动生成", "动态页面自动渲染\nCRUD 表单 + 数据表格\nAPI 端点即时可用", ACCENT2),
        ("04", "用户开始使用", "数据录入、搜索、导入导出\n关联管理、脚本校验\n全功能立即可用", PURPLE),
    ]

    x = Inches(0.5)
    for i, (num, title, desc, color) in enumerate(steps):
        add_rect(slide, x, Inches(1.8), Inches(2.8), Inches(3.0), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_flat_rect(slide, x, Inches(1.8), Inches(2.8), Inches(0.06), color)
        add_text(slide, x + Inches(0.2), Inches(2.0), Inches(0.6), Inches(0.5),
                 num, size=28, color=color, bold=True)
        add_text(slide, x + Inches(0.2), Inches(2.55), Inches(2.4), Inches(0.4),
                 title, size=15, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.2), Inches(3.05), Inches(2.4), Inches(1.5),
                 desc, size=11, color=LIGHT_TEXT)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           x + Inches(2.9), Inches(3.1), Inches(0.35), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PRIMARY
            arrow.line.fill.background()

        x += Inches(3.2)

    # 关键优势
    add_text(slide, Inches(0.8), Inches(5.1), Inches(6), Inches(0.4),
             "核心优势", size=18, color=WHITE, bold=True)

    advantages = [
        ("零代码扩展", "新增业务模块无需\n编写任何代码", PRIMARY),
        ("无迁移部署", "无需数据库迁移\nJSONB 自适应结构", ACCENT),
        ("实时生效", "配置修改后\n立即反映到 UI", ACCENT2),
        ("统一架构", "所有业务共享\n同一套 CRUD 引擎", PURPLE),
    ]

    x = Inches(0.8)
    for title, desc, color in advantages:
        add_rect(slide, x, Inches(5.5), Inches(2.8), Inches(1.5), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        dot = add_circle(slide, x + Inches(0.2), Inches(5.7), Inches(0.12), color)
        add_text(slide, x + Inches(0.45), Inches(5.65), Inches(2.2), Inches(0.3),
                 title, size=13, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.2), Inches(6.05), Inches(2.4), Inches(0.7),
                 desc, size=11, color=LIGHT_TEXT)
        x += Inches(3.05)


# ══════════════════════════════════════════════════════════
# 第 5 页 - 17 种字段控件
# ══════════════════════════════════════════════════════════
def slide_05_controls(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "17 种字段控件",
                     "覆盖企业数据管理中的各类字段需求，通过配置即可使用", 5)

    groups = [
        ("基础控件", PRIMARY, [
            ("Aa", "文本"), ("¶", "多行文本"), ("B", "富文本"),
            ("#", "数字"), ("D", "日期"), ("DT", "日期时间"),
        ]),
        ("选择控件", ACCENT, [
            ("▾", "下拉选择"), ("☑", "多选"), ("◉", "单选按钮"), ("☐", "复选框"),
        ]),
        ("文件控件", ACCENT2, [
            ("F", "文件上传"), ("I", "图片上传"),
        ]),
        ("关联控件", ACCENT3, [
            ("⇄", "关联(M:N)"), ("↓", "引用(父子)"), ("→", "引用选择"),
        ]),
        ("自动控件", PURPLE, [
            ("T", "自动时间戳"), ("S#", "自动编号"),
        ]),
    ]

    y = Inches(1.7)
    for group_name, color, controls in groups:
        # 组标签
        add_text(slide, Inches(0.8), y, Inches(1.5), Inches(0.35),
                 group_name, size=13, color=color, bold=True)

        # 控件卡片
        x = Inches(2.5)
        for icon, label in controls:
            add_rect(slide, x, y, Inches(1.6), Inches(0.7), CARD_BG,
                     border_color=BORDER, border_width=Pt(1))
            add_flat_rect(slide, x, y, Inches(0.05), Inches(0.7), color)
            add_text(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.4), Inches(0.3),
                     icon, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, x + Inches(0.5), y + Inches(0.18), Inches(1.0), Inches(0.3),
                     label, size=11, color=WHITE, bold=True)
            x += Inches(1.75)

        y += Inches(0.95)

    # 底部特性
    add_text(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
             "每种控件支持：必填校验 · 唯一性校验 · 正则表达式 · 默认值 · 选项数据源 · 可视化拖拽排序",
             size=12, color=LIGHT_TEXT)


# ══════════════════════════════════════════════════════════
# 第 6 页 - 数据关联体系
# ══════════════════════════════════════════════════════════
def slide_06_relations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "数据关联体系",
                     "三种灵活的关联模式，覆盖企业数据管理中的各类关系场景", 6)

    relations = [
        ("多对多关联", "Relation (M2M)", "⇄", PRIMARY,
         "双向关联，两端都可查看\n存储于 data_relations 表\n批量管理 + 导入自动匹配",
         "员工 ⇄ 项目\n标签 ⇄ 文章"),
        ("父子引用", "Reference (1:N)", "↓", ACCENT,
         "子记录存储父记录 ID\n支持字段继承 inheritFields\n父数据变更自动反映",
         "部门 → 员工\n类别 → 商品"),
        ("单向引用", "QuoteSelect", "→", ACCENT2,
         "单方向引用，ID 数组\n可多选目标记录\n不影响被引用方",
         "订单 → 商品\n计划 → 任务"),
    ]

    x = Inches(0.8)
    for title, subtitle, icon, color, desc, example in relations:
        add_rect(slide, x, Inches(1.7), Inches(3.7), Inches(4.0), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_flat_rect(slide, x, Inches(1.7), Inches(3.7), Inches(0.06), color)

        ic = add_circle(slide, x + Inches(1.35), Inches(2.0), Inches(0.9), color)
        add_text(slide, x + Inches(1.35), Inches(2.12), Inches(0.9), Inches(0.5),
                 icon, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.3), Inches(3.05), Inches(3.1), Inches(0.35),
                 title, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), Inches(3.4), Inches(3.1), Inches(0.25),
                 subtitle, size=11, color=color, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.3), Inches(3.85), Inches(3.1), Inches(1.2),
                 desc, size=11, color=LIGHT_TEXT)

        add_text(slide, x + Inches(0.3), Inches(5.0), Inches(0.8), Inches(0.25),
                 "示例:", size=10, color=color, bold=True)
        add_text(slide, x + Inches(1.0), Inches(5.0), Inches(2.4), Inches(0.5),
                 example, size=10, color=VERY_LIGHT)

        x += Inches(3.95)

    # 关系图谱
    add_rect(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.2), CARD_BG,
             border_color=BORDER, border_width=Pt(1))
    add_text(slide, Inches(1.1), Inches(6.0), Inches(3), Inches(0.4),
             "关系图谱可视化", size=16, color=WHITE, bold=True)
    add_text(slide, Inches(1.1), Inches(6.4), Inches(10), Inches(0.5),
             "基于 force-graph 力导向图展示数据关联网络  ·  单击展开节点  ·  双击跳转记录  ·  按集合名称自动着色  ·  三种关系类型不同线型",
             size=11, color=LIGHT_TEXT)


# ══════════════════════════════════════════════════════════
# 第 7 页 - 动态数据页面
# ══════════════════════════════════════════════════════════
def slide_07_dynamic_page(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "动态数据页面功能",
                     "一个 DynamicPage.vue 组件自动渲染所有业务数据集合的 CRUD 操作", 7)

    features = [
        ("🔍", "列表与搜索", "表格排序 · 分页浏览\n全字段关键词搜索\n列筛选 · 分页至1000条", PRIMARY),
        ("👁", "详情查看", "全字段格式化展示\n关联标签 · 图片预览\n引用跳转 · 图谱入口", ACCENT),
        ("✏", "创建 / 编辑", "动态表单 · 自动校验\n自动时间戳 · 自动编号\n乐观锁版本控制", ACCENT2),
        ("☑", "批量操作", "批量选择 · 批量删除\n仅管理员可批量操作\n引用约束检查", ACCENT3),
        ("📥", "数据导入", "Excel / JSON 模板导入\n自动字段映射\n关联名称自动解析", PRIMARY_DARK),
        ("📤", "数据导出", "标准 Excel 导出\n自定义脚本导出\nJSON/CSV/XML/TXT/HTML", PURPLE),
    ]

    x_start = Inches(0.5)
    y = Inches(1.7)

    for i, (icon, title, desc, color) in enumerate(features):
        col = i % 3
        row = i // 3
        x = x_start + col * Inches(4.2)
        cy = y + row * Inches(2.2)

        add_rect(slide, x, cy, Inches(3.9), Inches(2.0), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        ic = add_circle(slide, x + Inches(0.25), cy + Inches(0.25), Inches(0.55), color)
        add_text(slide, x + Inches(0.25), cy + Inches(0.32), Inches(0.55), Inches(0.35),
                 icon, size=16, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.95), cy + Inches(0.3), Inches(2.7), Inches(0.35),
                 title, size=15, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.25), cy + Inches(0.9), Inches(3.4), Inches(1.0),
                 desc, size=11, color=LIGHT_TEXT)

    # 底部亮点
    add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), CARD_BG,
             border_color=BORDER, border_width=Pt(1))
    add_text(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.7),
             "跨页导航:  点击关联标签 → 跳转目标记录   |   数据对比:  当前数据 vs 历史备份，字段级差异追踪   |   模板下载:  一键生成含字段说明的导入模板",
             size=11, color=LIGHT_TEXT)


# ══════════════════════════════════════════════════════════
# 第 8 页 - ETL 数据管道
# ══════════════════════════════════════════════════════════
def slide_08_etl(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "ETL 数据管道",
                     "可视化编排数据处理流程，6 种步骤类型灵活组合", 8)

    steps = [
        ("1", "HTTP 请求", "对接外部 API\n自定义请求头和请求体\n支持 GET/POST", PRIMARY),
        ("2", "JSON 输入", "手动输入数据\nJSON 数组或对象\n快速测试管道", ACCENT),
        ("3", "脚本转换", "Python 自定义处理\n接收 records 列表\n返回转换后数据", ACCENT2),
        ("4", "字段映射", "字段重命名\n源字段 → 目标字段\n可选保留未映射字段", ACCENT3),
        ("5", "数据过滤", "Python 条件表达式\n对每条记录求值\n筛选符合条件数据", PURPLE),
        ("6", "保存到集合", "写入目标数据集合\n自动关联处理\n完成 ETL 落地", PRIMARY_DARK),
    ]

    x = Inches(0.5)
    for i, (num, title, desc, color) in enumerate(steps):
        add_rect(slide, x, Inches(1.8), Inches(1.95), Inches(2.8), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_flat_rect(slide, x, Inches(1.8), Inches(1.95), Inches(0.05), color)
        nc = add_circle(slide, x + Inches(0.7), Inches(2.05), Inches(0.55), color)
        add_text(slide, x + Inches(0.7), Inches(2.13), Inches(0.55), Inches(0.35),
                 num, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), Inches(2.75), Inches(1.65), Inches(0.35),
                 title, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), Inches(3.15), Inches(1.65), Inches(1.2),
                 desc, size=10, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           x + Inches(2.0), Inches(3.0), Inches(0.2), Inches(0.15))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_TEXT
            arrow.line.fill.background()

        x += Inches(2.15)

    # 执行特性
    add_text(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(0.4),
             "执行与调试", size=18, color=WHITE, bold=True)

    exec_features = [
        ("试运行", "Dry-run 模式\n不写入数据库\n安全验证管道", PRIMARY),
        ("执行日志", "每步骤独立记录\n输入输出数据\n错误详细追踪", ACCENT),
        ("错误策略", "停止 / 跳过 / 继续\n灵活应对异常\n保障数据质量", ACCENT2),
        ("定时执行", "支持定时调度\n周期性数据同步\n自动化运维", PURPLE),
    ]

    x = Inches(0.8)
    for title, desc, color in exec_features:
        add_rect(slide, x, Inches(5.4), Inches(2.8), Inches(1.7), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        dot = add_circle(slide, x + Inches(0.2), Inches(5.6), Inches(0.12), color)
        add_text(slide, x + Inches(0.45), Inches(5.55), Inches(2.2), Inches(0.3),
                 title, size=13, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.2), Inches(5.95), Inches(2.4), Inches(0.9),
                 desc, size=11, color=LIGHT_TEXT)
        x += Inches(3.05)


# ══════════════════════════════════════════════════════════
# 第 9 页 - 脚本引擎
# ══════════════════════════════════════════════════════════
def slide_09_scripts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "脚本引擎",
                     "内置 Python 沙箱执行环境，支持校验脚本和导出脚本，通过 CodeMirror 6 在线编辑", 9)

    # 校验脚本
    add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.2), CARD_BG,
             border_color=PRIMARY, border_width=Pt(1.5))
    add_text(slide, Inches(1.1), Inches(1.85), Inches(4), Inches(0.4),
             "校验脚本 Validation Script", size=18, color=PRIMARY, bold=True)

    val_items = [
        "数据创建/编辑时自动执行校验",
        "返回错误列表阻止非法提交",
        "可返回警告信息（不阻止提交）",
        "可通过脚本自动设置关联关系",
        "按 PageConfig 绑定，每页一个脚本",
    ]
    y = Inches(2.4)
    for item in val_items:
        dot = add_circle(slide, Inches(1.3), y + Inches(0.05), Inches(0.1), PRIMARY)
        add_text(slide, Inches(1.55), y, Inches(4.5), Inches(0.3),
                 item, size=12, color=LIGHT_TEXT)
        y += Inches(0.4)

    # 校验脚本代码示例
    add_text(slide, Inches(1.1), Inches(4.6), Inches(3), Inches(0.3),
             "脚本示例:", size=12, color=PRIMARY, bold=True)
    code_bg = add_rect(slide, Inches(1.1), Inches(4.9), Inches(4.8), Inches(1.8),
                       RGBColor(0x0D, 0x13, 0x25))
    code = ("errors = []\nif not data.get('name'):\n"
            "    errors.append('名称不能为空')\n"
            "if data.get('score', 0) > 100:\n"
            "    errors.append('分数不能超过100')\n"
            "return errors, warnings")
    add_text(slide, Inches(1.3), Inches(5.0), Inches(4.5), Inches(1.6),
             code, size=10, color=ACCENT, font="Consolas")

    # 导出脚本
    add_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.2), CARD_BG,
             border_color=ACCENT2, border_width=Pt(1.5))
    add_text(slide, Inches(7.3), Inches(1.85), Inches(4), Inches(0.4),
             "导出脚本 Export Script", size=18, color=ACCENT2, bold=True)

    export_items = [
        "Python 脚本灵活定制导出格式",
        "支持 JSON / CSV / XML / TXT / HTML",
        "页面级导出：批量导出全部数据",
        "行级导出：单条记录定制导出",
        "按 PageConfig 绑定，可绑定多个",
    ]
    y = Inches(2.4)
    for item in export_items:
        dot = add_circle(slide, Inches(7.5), y + Inches(0.05), Inches(0.1), ACCENT2)
        add_text(slide, Inches(7.75), y, Inches(4.5), Inches(0.3),
                 item, size=12, color=LIGHT_TEXT)
        y += Inches(0.4)

    # 沙箱执行
    add_text(slide, Inches(7.3), Inches(4.6), Inches(3), Inches(0.3),
             "沙箱执行环境:", size=12, color=ACCENT2, bold=True)
    sandbox_items = [
        "隔离命名空间执行",
        "预注入 records / data / fields 变量",
        "内置 query() 函数查询其他集合",
        "CodeMirror 6 在线 Python 编辑器",
        "语法高亮 + 自动缩进",
    ]
    y = Inches(5.0)
    for item in sandbox_items:
        add_text(slide, Inches(7.5), y, Inches(4.5), Inches(0.3),
                 f"•  {item}", size=11, color=LIGHT_TEXT)
        y += Inches(0.35)


# ══════════════════════════════════════════════════════════
# 第 10 页 - 权限与安全
# ══════════════════════════════════════════════════════════
def slide_10_permissions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "权限与安全",
                     "基于角色的三级访问控制 (RBAC) + JWT 认证 + 操作审计，保障数据安全合规", 10)

    roles = [
        ("👑", "管理员 Admin", "完全访问", ACCENT2,
         ["系统管理（菜单/配置/用户）", "全部数据读写删除", "导入导出全功能",
          "数据对比与备份", "审计日志查看"]),
        ("🔧", "开发者 Developer", "配置读写", PRIMARY,
         ["访问已授权的菜单页面", "数据读写操作", "导入导出功能",
          "脚本/ETL 管理", "无法管理用户和备份"]),
        ("👤", "访客 Guest", "只读访问", LIGHT_TEXT,
         ["浏览已授权页面数据", "搜索和筛选", "下载导出文件",
          "无增删改权限", "无导入和批量操作"]),
    ]

    x = Inches(0.8)
    for icon, title, subtitle, color, perms in roles:
        add_rect(slide, x, Inches(1.7), Inches(3.7), Inches(3.8), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_flat_rect(slide, x, Inches(1.7), Inches(3.7), Inches(0.06), color)

        add_text(slide, x + Inches(0.3), Inches(1.95), Inches(0.5), Inches(0.4),
                 icon, size=22, color=color)
        add_text(slide, x + Inches(0.85), Inches(1.95), Inches(2.5), Inches(0.35),
                 title, size=15, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.85), Inches(2.3), Inches(2.5), Inches(0.25),
                 subtitle, size=11, color=color)

        y = Inches(2.75)
        for perm in perms:
            add_text(slide, x + Inches(0.3), y, Inches(3.1), Inches(0.28),
                     f"•  {perm}", size=11, color=LIGHT_TEXT)
            y += Inches(0.32)
        x += Inches(3.95)

    # 安全特性
    add_text(slide, Inches(0.8), Inches(5.8), Inches(6), Inches(0.4),
             "安全机制", size=18, color=WHITE, bold=True)

    sec_items = [
        ("JWT Token", "Bearer Token 认证\n自动过期机制", PRIMARY),
        ("菜单级权限", "按角色控制菜单可见性\n路由级访问保护", ACCENT),
        ("乐观锁", "_version 字段\n防止并发冲突覆盖", ACCENT2),
    ]

    x = Inches(0.8)
    for title, desc, color in sec_items:
        add_rect(slide, x, Inches(6.2), Inches(3.7), Inches(1.0), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        dot = add_circle(slide, x + Inches(0.2), Inches(6.35), Inches(0.1), color)
        add_text(slide, x + Inches(0.45), Inches(6.3), Inches(1.3), Inches(0.25),
                 title, size=12, color=WHITE, bold=True)
        add_text(slide, x + Inches(1.8), Inches(6.3), Inches(1.7), Inches(0.7),
                 desc, size=10, color=LIGHT_TEXT)
        x += Inches(3.95)


# ══════════════════════════════════════════════════════════
# 第 11 页 - Open API
# ══════════════════════════════════════════════════════════
def slide_11_open_api(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "Open API 外部接口",
                     "为外部系统提供安全的数据访问接口，API Key 认证 + 按集合粒度开放", 11)

    # 工作流程
    flow = [
        ("1", "管理员生成 API Key", "在 API 密钥管理页面\n生成 cm_ 前缀的密钥", PRIMARY),
        ("2", "配置公开集合", "在 PageConfig 中\n开启 Open API 开关", ACCENT),
        ("3", "外部系统调用", "携带 X-API-Key 头\n访问 RESTful 接口", ACCENT2),
    ]

    x = Inches(0.8)
    for num, title, desc, color in flow:
        add_rect(slide, x, Inches(1.8), Inches(3.7), Inches(1.8), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        nc = add_circle(slide, x + Inches(0.2), Inches(2.0), Inches(0.5), color)
        add_text(slide, x + Inches(0.2), Inches(2.08), Inches(0.5), Inches(0.35),
                 num, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.85), Inches(2.0), Inches(2.5), Inches(0.35),
                 title, size=14, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.85), Inches(2.4), Inches(2.6), Inches(0.8),
                 desc, size=11, color=LIGHT_TEXT)
        x += Inches(3.95)

    # API 特性
    add_text(slide, Inches(0.8), Inches(3.9), Inches(4), Inches(0.4),
             "接口规范", size=18, color=WHITE, bold=True)

    # 代码示例
    add_rect(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.8), RGBColor(0x0D, 0x13, 0x25),
             border_color=BORDER, border_width=Pt(1))
    code = ("# 请求示例\n"
            "GET /open-api/inspection-case\n"
            "Headers:\n"
            "  X-API-Key: cm_xxxxxx\n"
            "\n"
            "# 响应格式\n"
            "[\n"
            '  { "id": "rec-001",\n'
            '    "data": { "name": "用例1" },\n'
            '    "created_at": "2024-01-01" }\n'
            "]")
    add_text(slide, Inches(1.0), Inches(4.4), Inches(5.1), Inches(2.6),
             code, size=10, color=ACCENT, font="Consolas")

    # 右侧特性列表
    add_rect(slide, Inches(7.0), Inches(4.3), Inches(5.5), Inches(2.8), CARD_BG,
             border_color=BORDER, border_width=Pt(1))
    add_text(slide, Inches(7.3), Inches(4.4), Inches(3), Inches(0.35),
             "API 特性", size=15, color=WHITE, bold=True)

    api_items = [
        ("认证方式", "API Key (X-API-Key 请求头)", PRIMARY),
        ("访问粒度", "按集合开放，每个集合独立控制", ACCENT),
        ("数据格式", "RESTful JSON 响应", ACCENT2),
        ("密钥管理", "管理员可随时生成/撤销密钥", ACCENT3),
        ("安全隔离", "仅可访问标记为公开的集合", PURPLE),
    ]
    y = Inches(4.9)
    for title, desc, color in api_items:
        dot = add_circle(slide, Inches(7.5), y + Inches(0.05), Inches(0.1), color)
        add_text(slide, Inches(7.75), y, Inches(1.6), Inches(0.28),
                 title, size=12, color=WHITE, bold=True)
        add_text(slide, Inches(9.4), y, Inches(2.8), Inches(0.28),
                 desc, size=11, color=LIGHT_TEXT)
        y += Inches(0.4)


# ══════════════════════════════════════════════════════════
# 第 12 页 - 备份与恢复
# ══════════════════════════════════════════════════════════
def slide_12_backup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "备份与恢复",
                     "手动/自动备份 + 跨环境迁移 + 数据对比，保障数据安全", 12)

    # 备份功能
    add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3.0), CARD_BG,
             border_color=PRIMARY, border_width=Pt(1.5))
    add_text(slide, Inches(1.1), Inches(1.85), Inches(4), Inches(0.4),
             "备份管理", size=18, color=PRIMARY, bold=True)

    backup_items = [
        "手动一键备份当前数据库快照",
        "自动定时备份：小时 / 天 / 周 / 月",
        "保留策略：可配置保留数量",
        "下载备份 .zip 文件到本地",
        "上传外部备份文件用于迁移",
        "一键恢复至任意历史版本",
    ]
    y = Inches(2.4)
    for item in backup_items:
        dot = add_circle(slide, Inches(1.3), y + Inches(0.05), Inches(0.1), PRIMARY)
        add_text(slide, Inches(1.55), y, Inches(4.5), Inches(0.3),
                 item, size=12, color=LIGHT_TEXT)
        y += Inches(0.37)

    # 数据对比
    add_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(3.0), CARD_BG,
             border_color=ACCENT, border_width=Pt(1.5))
    add_text(slide, Inches(7.3), Inches(1.85), Inches(4), Inches(0.4),
             "数据对比", size=18, color=ACCENT, bold=True)

    compare_items = [
        "当前数据 vs 任意历史备份",
        "两个备份版本之间的对比",
        "字段级差异追踪（旧值 → 新值）",
        "新增 / 修改 / 删除 分类标记",
        "变更字段高亮显示",
        "对比结果导出 Excel 报告",
    ]
    y = Inches(2.4)
    for item in compare_items:
        dot = add_circle(slide, Inches(7.5), y + Inches(0.05), Inches(0.1), ACCENT)
        add_text(slide, Inches(7.75), y, Inches(4.5), Inches(0.3),
                 item, size=12, color=LIGHT_TEXT)
        y += Inches(0.37)

    # 备份调度示意
    add_text(slide, Inches(0.8), Inches(5.0), Inches(6), Inches(0.4),
             "自动备份调度", size=18, color=WHITE, bold=True)

    schedules = [
        ("每小时", "hourly", "高频变更\n实时数据安全", PRIMARY),
        ("每天", "daily", "日常运营\n推荐默认方案", ACCENT),
        ("每周", "weekly", "中等频率\n节省存储空间", ACCENT2),
        ("每月", "monthly", "长期归档\n合规审计需求", PURPLE),
    ]

    x = Inches(0.8)
    for label, freq, desc, color in schedules:
        add_rect(slide, x, Inches(5.4), Inches(2.8), Inches(1.7), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_flat_rect(slide, x, Inches(5.4), Inches(2.8), Inches(0.05), color)
        add_text(slide, x + Inches(0.2), Inches(5.55), Inches(2.4), Inches(0.3),
                 label, size=15, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.2), Inches(5.85), Inches(2.4), Inches(0.3),
                 freq, size=10, color=color, font="Consolas")
        add_text(slide, x + Inches(0.2), Inches(6.2), Inches(2.4), Inches(0.6),
                 desc, size=11, color=LIGHT_TEXT)
        x += Inches(3.05)


# ══════════════════════════════════════════════════════════
# 第 13 页 - 操作审计
# ══════════════════════════════════════════════════════════
def slide_13_audit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "操作审计日志",
                     "完整记录所有数据变更操作，多维度查询与分析，满足合规审计需求", 13)

    # 记录的操作类型
    add_text(slide, Inches(0.8), Inches(1.7), Inches(4), Inches(0.4),
             "记录的操作类型", size=18, color=WHITE, bold=True)

    ops = [
        ("创建", "新增记录", PRIMARY),
        ("更新", "编辑记录", ACCENT),
        ("删除", "删除记录", ACCENT3),
        ("导入", "批量导入", ACCENT2),
        ("导出", "数据导出", PURPLE),
        ("备份", "备份/恢复", PRIMARY_DARK),
    ]

    x = Inches(0.8)
    for op, desc, color in ops:
        add_rect(slide, x, Inches(2.1), Inches(1.85), Inches(0.9), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_flat_rect(slide, x, Inches(2.1), Inches(0.05), Inches(0.9), color)
        add_text(slide, x + Inches(0.15), Inches(2.2), Inches(1.5), Inches(0.3),
                 op, size=14, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.15), Inches(2.5), Inches(1.5), Inches(0.3),
                 desc, size=10, color=LIGHT_TEXT)
        x += Inches(2.05)

    # 查询能力
    add_text(slide, Inches(0.8), Inches(3.3), Inches(6), Inches(0.4),
             "多维度查询", size=18, color=WHITE, bold=True)

    filters = [
        ("日期范围", "按时间段筛选操作记录\n快速定位特定时间的变更"),
        ("操作类型", "按创建/更新/删除/导入等类型\n筛选特定类型的操作"),
        ("目标类型", "按数据集合（业务模块）\n查看特定模块的操作历史"),
        ("操作人", "按执行人筛选\n追踪特定用户的操作行为"),
    ]

    x = Inches(0.8)
    for title, desc in filters:
        add_rect(slide, x, Inches(3.7), Inches(2.8), Inches(1.5), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_text(slide, x + Inches(0.2), Inches(3.85), Inches(2.4), Inches(0.3),
                 title, size=13, color=PRIMARY, bold=True)
        add_text(slide, x + Inches(0.2), Inches(4.2), Inches(2.4), Inches(0.8),
                 desc, size=11, color=LIGHT_TEXT)
        x += Inches(3.05)

    # 特色功能
    add_text(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.4),
             "特色功能", size=18, color=WHITE, bold=True)

    specials = [
        ("批次分组", "关联操作自动聚合\n如：批量导入的所有记录\n统一 Batch ID 标识", ACCENT),
        ("变更详情", "记录每个字段的变更\n旧值 → 新值 对比\n字段标签 + 显示名", ACCENT2),
        ("导出报告", "筛选结果导出 Excel\n完整操作轨迹\n满足合规审计要求", PURPLE),
    ]

    x = Inches(0.8)
    for title, desc, color in specials:
        add_rect(slide, x, Inches(5.9), Inches(3.7), Inches(1.4), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        dot = add_circle(slide, x + Inches(0.2), Inches(6.05), Inches(0.12), color)
        add_text(slide, x + Inches(0.45), Inches(6.0), Inches(3.0), Inches(0.3),
                 title, size=13, color=WHITE, bold=True)
        add_text(slide, x + Inches(0.2), Inches(6.35), Inches(3.3), Inches(0.8),
                 desc, size=10, color=LIGHT_TEXT)
        x += Inches(3.95)


# ══════════════════════════════════════════════════════════
# 第 14 页 - 系统管理
# ══════════════════════════════════════════════════════════
def slide_14_admin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_slide_header(slide, "系统管理",
                     "10 个管理页面，覆盖系统配置、数据工具、运维管理全流程", 14)

    pages = [
        ("菜", "菜单管理", "可视化树形编辑\n1-3 级层级菜单\n角色可见性控制", PRIMARY),
        ("配", "页面配置", "字段架构定义\n17 种控件类型\n绑定脚本/API", ACCENT),
        ("人", "用户管理", "三级角色体系\n账号增删改\n密码重置", ACCENT2),
        ("E", "导出脚本", "Python 自定义导出\n5 种输出格式\n页面级/行级", ACCENT3),
        ("V", "校验脚本", "数据校验规则\n阻止非法提交\n返回警告信息", PURPLE),
        ("ETL", "ETL 任务", "数据管道编排\n6 种步骤类型\n试运行/执行日志", PRIMARY_DARK),
        ("K", "API 密钥", "Open API 管理\n密钥生成/撤销\n外部系统接入", RGBColor(0xE0, 0x6C, 0x9F)),
        ("日", "操作日志", "全操作审计\n多维度查询\n导出报告", RGBColor(0x6C, 0xB4, 0xE0)),
        ("备", "备份管理", "手动/自动备份\n恢复/下载/上传\n定时调度", RGBColor(0xE0, 0xC0, 0x6C)),
        ("对", "数据对比", "备份快照对比\n字段级差异\n变更高亮", RGBColor(0x6C, 0xE0, 0xA0)),
    ]

    x_start = Inches(0.5)
    y_start = Inches(1.7)

    for i, (icon, title, desc, color) in enumerate(pages):
        col = i % 5
        row = i // 5
        x = x_start + col * Inches(2.5)
        y = y_start + row * Inches(2.8)

        add_rect(slide, x, y, Inches(2.3), Inches(2.5), CARD_BG,
                 border_color=BORDER, border_width=Pt(1))
        add_flat_rect(slide, x, y, Inches(2.3), Inches(0.05), color)

        ic = add_circle(slide, x + Inches(0.8), y + Inches(0.25), Inches(0.6), color)
        add_text(slide, x + Inches(0.8), y + Inches(0.33), Inches(0.6), Inches(0.4),
                 icon, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

        add_text(slide, x + Inches(0.15), y + Inches(1.0), Inches(2.0), Inches(0.35),
                 title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), y + Inches(1.4), Inches(2.0), Inches(1.0),
                 desc, size=10, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# 第 15 页 - 总结
# ══════════════════════════════════════════════════════════
def slide_15_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_circle(slide, Inches(10.0), Inches(-1.5), Inches(4.5), RGBColor(0x1E, 0x3A, 0x6E))
    add_circle(slide, Inches(11.5), Inches(0.0), Inches(3.0), RGBColor(0x24, 0x44, 0x7A))
    add_circle(slide, Inches(-1.5), Inches(5.0), Inches(3.5), RGBColor(0x1E, 0x3A, 0x6E))

    add_flat_rect(slide, 0, 0, Inches(0.12), SLIDE_HEIGHT, PRIMARY)

    add_text(slide, Inches(1.2), Inches(0.8), Inches(8), Inches(0.8),
             "平台核心优势", size=36, color=WHITE, bold=True)

    add_flat_rect(slide, Inches(1.2), Inches(1.6), Inches(1.5), Inches(0.04), PRIMARY)

    advantages = [
        ("配置驱动", "零代码扩展新业务模块，无需数据库迁移", PRIMARY),
        ("17 种控件", "覆盖文本、选择、文件、关联、自动字段等全场景", ACCENT),
        ("三种关联", "M:N 双向 + 父子引用 + 单向引用，灵活数据建模", ACCENT2),
        ("ETL 管道", "6 种步骤类型可视化编排，试运行 + 执行日志", ACCENT3),
        ("脚本引擎", "Python 沙箱执行校验与导出脚本，CodeMirror 在线编辑", PURPLE),
        ("权限体系", "三级 RBAC + JWT 认证 + 菜单级权限 + 乐观锁", PRIMARY_DARK),
        ("审计追踪", "全操作记录、批次分组、多维查询、导出报告", RGBColor(0x6C, 0xB4, 0xE0)),
        ("备份恢复", "手动/自动备份、跨环境迁移、字段级数据对比", RGBColor(0xE0, 0xC0, 0x6C)),
    ]

    y = Inches(1.9)
    for title, desc, color in advantages:
        dot = add_circle(slide, Inches(1.4), y + Inches(0.08), Inches(0.14), color)
        add_text(slide, Inches(1.7), y, Inches(2.2), Inches(0.35),
                 title, size=15, color=WHITE, bold=True)
        add_text(slide, Inches(4.0), y, Inches(7.5), Inches(0.35),
                 desc, size=13, color=LIGHT_TEXT)
        y += Inches(0.5)

    # 测试覆盖
    add_flat_rect(slide, Inches(1.2), Inches(6.0), Inches(11), Inches(0.04),
                  RGBColor(0x30, 0x3D, 0x5F))

    add_text(slide, Inches(1.2), Inches(6.2), Inches(11), Inches(0.5),
             "质量保障：326 前端测试 (Vitest) + 241 后端测试 (Pytest) = 567 单元测试全覆盖",
             size=14, color=PRIMARY)

    # 技术标签
    tags = ["Vue 3", "TypeScript", "Element Plus", "Flask", "PostgreSQL", "force-graph"]
    x = Inches(1.2)
    for tag in tags:
        w = add_tag(slide, x, Inches(6.8), tag)
        x += w + Inches(0.15)


# ══════════════════════════════════════════════════════════
# 主函数
# ══════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_01_cover(prs)           # 1. 封面
    slide_02_overview(prs)        # 2. 平台概述
    slide_03_architecture(prs)    # 3. 系统架构
    slide_04_config_driven(prs)   # 4. 配置驱动
    slide_05_controls(prs)        # 5. 17 种字段控件
    slide_06_relations(prs)       # 6. 数据关联体系
    slide_07_dynamic_page(prs)    # 7. 动态数据页面
    slide_08_etl(prs)             # 8. ETL 数据管道
    slide_09_scripts(prs)         # 9. 脚本引擎
    slide_10_permissions(prs)     # 10. 权限与安全
    slide_11_open_api(prs)        # 11. Open API
    slide_12_backup(prs)          # 12. 备份与恢复
    slide_13_audit(prs)           # 13. 操作审计
    slide_14_admin(prs)           # 14. 系统管理
    slide_15_summary(prs)         # 15. 总结

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "动态数据管理平台-系统功能介绍-v2.pptx")
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")


if __name__ == "__main__":
    main()
