#!/usr/bin/env python3
"""Generate a technical introduction PPT for check-manage system."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
BG_DARK   = RGBColor(0x1B, 0x1F, 0x3B)
BG_MEDIUM = RGBColor(0x24, 0x2B, 0x4E)
BG_CARD   = RGBColor(0x2D, 0x35, 0x5E)
ACCENT    = RGBColor(0x64, 0xB5, 0xF6)
ACCENT2   = RGBColor(0x4D, 0xD0, 0xE1)
GREEN     = RGBColor(0x66, 0xBB, 0x6A)
ORANGE    = RGBColor(0xFF, 0xA7, 0x26)
RED       = RGBColor(0xEF, 0x53, 0x50)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xB0, 0xBE, 0xC5)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
CODE_BG   = RGBColor(0x26, 0x32, 0x38)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p


def add_para(tf, text, size=14, color=WHITE, bold=False, space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p


def add_bullet(tf, text, size=14, color=WHITE, level=0, bold=False, font_name='Microsoft YaHei'):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(2)
    run = p.add_run()
    prefix = '  ' * level + ('• ' if level == 0 else '◦ ')
    run.text = prefix + text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return p


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = add_shape(slide, left, top, width, height, CODE_BG, corner_radius=Inches(0.1))
    tf = shape.text_frame
    tf.word_wrap = True
    set_text(tf, code_text, size=font_size, color=ACCENT2, font_name='Consolas')
    return shape


def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    # Title
    tb = add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5))
    set_text(tb.text_frame, 'check-manage 系统技术介绍', size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Subtitle
    tb2 = add_text_box(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.8))
    set_text(tb2.text_frame, '配置驱动的动态数据管理平台', size=24, color=ACCENT, alignment=PP_ALIGN.CENTER)
    # Description
    tb3 = add_text_box(slide, Inches(2), Inches(4.8), Inches(9), Inches(1.5))
    tf = tb3.text_frame
    tf.word_wrap = True
    set_text(tf, 'Vue 3 + Flask + PostgreSQL JSONB', size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_para(tf, '面向软件工程师的系统架构与功能详解', size=16, color=GRAY, alignment=PP_ALIGN.CENTER)


def make_toc_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    set_text(tb.text_frame, '目录', size=32, color=WHITE, bold=True)

    items = [
        ('01', '系统架构总览', '技术栈、数据库设计、核心理念'),
        ('02', '配置驱动机制', 'PageConfig → 动态 UI → 通用 API'),
        ('03', '16 种字段控件', '从文本到关联选择，全类型详解'),
        ('04', '数据关联原理', 'M:N / Reference / QuoteSelect 三种关联'),
        ('05', '乐观锁与冲突解决', 'version 字段 + 双重校验机制'),
        ('06', '查询引擎', 'MongoDB 语法 → PostgreSQL JSONB 翻译'),
        ('07', '校验与导出脚本', 'Python 沙箱 + 自定义数据处理'),
        ('08', 'ETL 数据管道', 'HTTP 抽取 → 脚本转换 → 字段映射'),
        ('09', 'Open API', 'API Key 认证 + 外部系统集成'),
        ('10', '权限与审计', 'RBAC 三级角色 + 全量操作日志'),
        ('11', '备份与恢复', '手动/定时备份 + 数据对比 + 跨环境迁移'),
        ('12', '部署方案', '本地部署 + frp 公网穿透'),
    ]
    y = Inches(1.4)
    for idx, (num, title, desc) in enumerate(items):
        col = idx % 3
        row = idx // 3
        x = Inches(0.6 + col * 4.1)
        cy = y + Inches(row * 1.35)
        card = add_shape(slide, x, cy, Inches(3.8), Inches(1.15), BG_CARD)
        tf = card.text_frame
        tf.word_wrap = True
        set_text(tf, f'  {num}  {title}', size=15, color=ACCENT, bold=True)
        add_para(tf, f'       {desc}', size=11, color=GRAY)


def make_arch_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '01  系统架构总览', size=28, color=WHITE, bold=True)

    # Left: tech stack
    card1 = add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.8), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  技术架构', size=18, color=ACCENT, bold=True)
    lines = [
        '',
        '  ┌─────────────── 浏览器 ───────────────┐',
        '  │  Vue 3 + Element Plus + Pinia        │',
        '  │  vue-codemirror + ECharts + XLSX      │',
        '  └──────────────┬────────────────────────┘',
        '                 │  HTTP / JWT Bearer',
        '  ┌──────────────▼────────────────────────┐',
        '  │  Flask 后端 (port 3001)               │',
        '  │  Blueprint 路由 × 12 个模块            │',
        '  │  反向代理 (port 8080, 生产环境)        │',
        '  └──────────────┬────────────────────────┘',
        '                 │  psycopg2',
        '  ┌──────────────▼────────────────────────┐',
        '  │  PostgreSQL                           │',
        '  │  13 张表 + JSONB 灵活存储              │',
        '  └───────────────────────────────────────┘',
    ]
    for line in lines:
        add_para(tf, line, size=12, color=ACCENT2, font_name='Consolas')

    # Right: database design
    card2 = add_shape(slide, Inches(6.7), Inches(1.3), Inches(6.1), Inches(5.8), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  核心数据库表', size=18, color=ACCENT, bold=True)
    tables = [
        ('dynamic_data',   'JSONB 单表存储所有业务数据'),
        ('page_configs',   '页面字段配置（驱动 UI 渲染）'),
        ('data_relations',  'M:N 关联中间表'),
        ('menus',          '树形菜单 + RBAC 角色'),
        ('users',          '用户表（admin/developer/guest）'),
        ('operation_logs',  '全量操作审计日志'),
        ('backups',        '备份记录与元数据'),
        ('export_scripts',  '自定义导出脚本'),
        ('validation_scripts', '数据校验脚本'),
        ('etl_tasks / etl_logs', 'ETL 任务与执行日志'),
        ('api_keys',        'Open API 密钥管理'),
    ]
    for name, desc in tables:
        add_para(tf, f'  {name}', size=13, color=GREEN, bold=True, font_name='Consolas')
        add_para(tf, f'     {desc}', size=11, color=GRAY, space_before=Pt(0))

    # Bottom note
    tb2 = add_text_box(slide, Inches(0.5), Inches(7.15), Inches(12), Inches(0.3))
    set_text(tb2.text_frame, '核心理念：所有业务数据存储在 dynamic_data 一张表，通过 collection 字段区分实体，JSONB data 列存储字段值', size=11, color=GRAY)


def make_config_driven_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '02  配置驱动机制', size=28, color=WHITE, bold=True)

    # Flow
    card = add_shape(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.2), BG_CARD)
    tf = card.text_frame
    tf.word_wrap = True
    set_text(tf, '  从配置到运行的完整流程', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=6, color=WHITE)
    flow = '  page_configs.fields (JSON)  →  前端 DynamicPage.vue 读取  →  DynamicForm 渲染表单  →  DataTable 渲染列表'
    add_para(tf, flow, size=13, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  后端 dynamic.py 通用 CRUD  →  不需要为新实体写任何代码  →  新增业务实体只需 INSERT 两行 SQL', size=13, color=GREEN, font_name='Consolas')

    # Left: how to add entity
    card2 = add_shape(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(3.5), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  新增业务实体（零代码）', size=16, color=ACCENT, bold=True)
    code = '''
  -- 1. 定义字段配置
  INSERT INTO page_configs (id, name, fields)
  VALUES ('page-products', '产品管理',
    '[{"fieldName":"name",
       "label":"产品名称",
       "controlType":"text",
       "required":true}]');

  -- 2. 创建菜单入口
  INSERT INTO menus (id, name, path, page_id)
  VALUES ('menu-prod', '产品', '/products',
          'page-products');

  -- 完成！自动获得 UI + API'''
    add_para(tf, code, size=11, color=ACCENT2, font_name='Consolas')

    # Right: what you get
    card3 = add_shape(slide, Inches(6.8), Inches(3.8), Inches(6), Inches(3.5), BG_CARD)
    tf = card3.text_frame
    tf.word_wrap = True
    set_text(tf, '  自动获得的能力', size=16, color=ACCENT, bold=True)
    capabilities = [
        ('前端页面', '表单 + 表格 + 搜索 + 分页'),
        ('REST API', 'GET/POST/PUT/DELETE /<collection>'),
        ('数据校验', '必填 + 类型 + 自定义脚本'),
        ('导入导出', 'Excel 批量导入 + 自定义导出'),
        ('操作审计', '自动记录增删改操作日志'),
        ('关联关系', 'relation / reference / quoteSelect'),
        ('Open API', '外部系统 API Key 访问'),
        ('MongoDB 查询', '查询控制台支持复杂查询'),
    ]
    for name, desc in capabilities:
        add_para(tf, f'  ✓ {name}', size=13, color=GREEN, bold=True)
        add_para(tf, f'     {desc}', size=11, color=GRAY, space_before=Pt(0))


def make_controls_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '03  16 种字段控件', size=28, color=WHITE, bold=True)

    controls = [
        ('text',           'TextInput',         '单行文本'),
        ('textarea',       'TextArea',          '多行文本'),
        ('number',         'NumberInput',       '数字输入'),
        ('select',         'SelectInput',       '下拉单选'),
        ('multiSelect',    'MultiSelect',       '多选标签'),
        ('radio',          'RadioGroup',        '单选按钮组'),
        ('checkbox',       'CheckboxGroup',     '多选框组'),
        ('date',           'DatePicker',        '日期选择'),
        ('datetime',       'DatePicker',        '日期时间'),
        ('file',           'FileUpload',        '文件上传'),
        ('image',          'ImageUpload',       '图片上传'),
        ('richText',       'RichTextEditor',    '富文本编辑'),
        ('autoSequence',   'AutoSequence',      '自增序列'),
        ('autoTimestamp',  'AutoTimestamp',      '自动时间戳'),
        ('reference',      'ReferenceSelect',   '父子引用'),
        ('relation',       'RelationSelect',    'M:N 多对多'),
    ]

    for idx, (ctrl_type, component, desc) in enumerate(controls):
        col = idx % 4
        row = idx // 4
        x = Inches(0.4 + col * 3.2)
        y = Inches(1.3 + row * 1.5)
        card = add_shape(slide, x, y, Inches(3.0), Inches(1.3), BG_CARD)
        tf = card.text_frame
        tf.word_wrap = True
        set_text(tf, f'  {ctrl_type}', size=14, color=ACCENT, bold=True, font_name='Consolas')
        add_para(tf, f'  {component}.vue', size=10, color=GRAY, font_name='Consolas')
        add_para(tf, f'  {desc}', size=12, color=WHITE)

    # QuoteSelect note
    tb2 = add_text_box(slide, Inches(0.5), Inches(7.15), Inches(12), Inches(0.3))
    set_text(tb2.text_frame, '另有 quoteSelect（引用选择）控件，controlType 为 quoteSelect，用于从关联集合中选择记录并存储 ID 数组', size=11, color=GRAY)


def make_relation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '04  数据关联原理', size=28, color=WHITE, bold=True)

    # Type 1: M:N relation
    card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(4.0), Inches(5.9), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  relation (M:N 多对多)', size=15, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  存储方式', size=12, color=ORANGE, bold=True)
    add_para(tf, '  通过 data_relations 中间表', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  中间表结构', size=12, color=ORANGE, bold=True)
    code1 = '''  collection | record_id
  field_name | related_collection
  related_id'''
    add_para(tf, code1, size=10, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  双向同步', size=12, color=ORANGE, bold=True)
    add_para(tf, '  巡检用例 ←→ 巡检模板', size=11, color=WHITE)
    add_para(tf, '  A 关联 B 时，自动在 B 侧', size=11, color=WHITE)
    add_para(tf, '  创建反向关联', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  删除保护', size=12, color=RED, bold=True)
    add_para(tf, '  被关联的记录不可删除', size=11, color=WHITE)

    # Type 2: reference
    card2 = add_shape(slide, Inches(4.7), Inches(1.2), Inches(4.0), Inches(5.9), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  reference (父子引用)', size=15, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  存储方式', size=12, color=ORANGE, bold=True)
    add_para(tf, '  data JSONB 中直接存父记录 ID', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  数据结构', size=12, color=ORANGE, bold=True)
    code2 = '''  // PPT配置记录
  {
    "用例ID": "parent-record-id",
    "操作类型": "检查"
  }'''
    add_para(tf, code2, size=10, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  字段继承', size=12, color=ORANGE, bold=True)
    add_para(tf, '  子记录自动继承父记录', size=11, color=WHITE)
    add_para(tf, '  的指定字段值', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  级联删除保护', size=12, color=RED, bold=True)
    add_para(tf, '  有子记录时父记录不可删', size=11, color=WHITE)

    # Type 3: quoteSelect
    card3 = add_shape(slide, Inches(8.9), Inches(1.2), Inches(4.0), Inches(5.9), BG_CARD)
    tf = card3.text_frame
    tf.word_wrap = True
    set_text(tf, '  quoteSelect (引用选择)', size=15, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  存储方式', size=12, color=ORANGE, bold=True)
    add_para(tf, '  data JSONB 中存 ID 数组', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  数据结构', size=12, color=ORANGE, bold=True)
    code3 = '''  // PPT配置记录
  {
    "模板ID": [
      "template-001",
      "template-002"
    ]
  }'''
    add_para(tf, code3, size=10, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  使用场景', size=12, color=ORANGE, bold=True)
    add_para(tf, '  从目标集合中选择多条记录', size=11, color=WHITE)
    add_para(tf, '  轻量级引用，无中间表', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  适合', size=12, color=GREEN, bold=True)
    add_para(tf, '  单向引用、不需要反查', size=11, color=WHITE)


def make_conflict_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '05  乐观锁与冲突解决', size=28, color=WHITE, bold=True)

    # Left: principle
    card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.9), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  乐观锁机制 (Optimistic Locking)', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=6, color=WHITE)
    add_para(tf, '  原理', size=14, color=ORANGE, bold=True)
    add_para(tf, '  每条记录维护 version 字段（整数递增）', size=12, color=WHITE)
    add_para(tf, '  更新时携带 version，服务端比对后决定是否写入', size=12, color=WHITE)
    add_para(tf, '', size=6, color=WHITE)
    add_para(tf, '  双重校验', size=14, color=ORANGE, bold=True)
    add_para(tf, '  第 1 层：SELECT 查出 db_version，与 client_version 比对', size=11, color=WHITE)
    add_para(tf, '  第 2 层：UPDATE ... WHERE version = db_version', size=11, color=WHITE)
    add_para(tf, '  如果 rowcount == 0，说明 SELECT 和 UPDATE 之间又被修改', size=11, color=WHITE)
    add_para(tf, '', size=6, color=WHITE)
    add_para(tf, '  冲突时返回', size=14, color=RED, bold=True)
    code = '''  HTTP 409 Conflict
  {
    "error": "数据已被其他用户修改，请刷新后重试",
    "code": "VERSION_CONFLICT",
    "_version": <current_db_version>
  }'''
    add_para(tf, code, size=10, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=6, color=WHITE)
    add_para(tf, '  前端处理', size=14, color=GREEN, bold=True)
    add_para(tf, '  pageConfig.ts 更新时自动携带 _version', size=11, color=WHITE)
    add_para(tf, '  收到 409 后提示用户刷新重试', size=11, color=WHITE)

    # Right: flow
    card2 = add_shape(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.9), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  时序流程', size=16, color=ACCENT, bold=True)
    flow_lines = [
        '',
        '  用户A 读取记录 (version=3)',
        '  用户B 读取记录 (version=3)',
        '  ',
        '  ── 用户A 先提交 ──────────────',
        '  PUT {_version: 3, data...}',
        '  → SELECT: db_version=3 ✓',
        '  → UPDATE WHERE version=3',
        '  → version 变为 4',
        '  → 返回 200 {_version: 4}',
        '  ',
        '  ── 用户B 后提交 ──────────────',
        '  PUT {_version: 3, data...}',
        '  → SELECT: db_version=4 ≠ 3',
        '  → 返回 409 VERSION_CONFLICT',
        '  ',
        '  ── 用户B 刷新后重试 ──────────',
        '  GET → 获取最新数据 (version=4)',
        '  修改后 PUT {_version: 4}',
        '  → 成功',
    ]
    for line in flow_lines:
        color = RED if '409' in line or '≠' in line else GREEN if '✓' in line or '成功' in line else ACCENT2
        add_para(tf, line, size=11, color=color, font_name='Consolas')


def make_query_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '06  查询引擎', size=28, color=WHITE, bold=True)

    # Left: syntax
    card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.9), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  MongoDB 风格查询语法', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  utils/mongo_query.py 将 MongoDB 语法翻译为 PostgreSQL', size=11, color=GRAY)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  支持 12 种操作符', size=13, color=ORANGE, bold=True)
    ops = [
        '比较:  $eq $ne $gt $gte $lt $lte',
        '数组:  $in $nin',
        '字符串: $regex $like',
        '元素:  $exists $size',
        '逻辑:  $and $or $not $nor',
    ]
    for op in ops:
        add_para(tf, f'  {op}', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  查询示例', size=13, color=ORANGE, bold=True)
    code = '''  {
    "collection": "inspection-case",
    "query": {
      "$or": [
        {"用例类型": "功能巡检"},
        {"优先级": "高"}
      ],
      "用例名称": {"$like": "登录"}
    },
    "sort": {"用例ID": 1},
    "limit": 200
  }'''
    add_para(tf, code, size=10, color=ACCENT2, font_name='Consolas')

    # Right: features
    card2 = add_shape(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(2.6), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  中文标签自动映射', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  查询中使用中文标签（如"用例ID"）', size=12, color=WHITE)
    add_para(tf, '  系统通过 page_configs.fields 自动转换', size=12, color=WHITE)
    add_para(tf, '  为实际字段名（如"caseid"）', size=12, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  remap_labels(query, fields)', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  → 递归遍历 query 对象中的所有 key', size=11, color=GRAY)

    card3 = add_shape(slide, Inches(6.9), Inches(4.1), Inches(6.0), Inches(3.0), BG_CARD)
    tf = card3.text_frame
    tf.word_wrap = True
    set_text(tf, '  连表查询 (lookup)', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  支持三种关联类型的跨集合 JOIN', size=12, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  relation → 通过 data_relations 表', size=11, color=GREEN, font_name='Consolas')
    add_para(tf, '  reference → 本地 ID 字段直接关联', size=11, color=GREEN, font_name='Consolas')
    add_para(tf, '  quoteSelect → 本地 ID 数组展开', size=11, color=GREEN, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  查询控制台', size=13, color=ORANGE, bold=True)
    add_para(tf, '  CodeMirror 编辑器 + 字段自动补全', size=11, color=WHITE)
    add_para(tf, '  实时结果展示 + Excel 导出', size=11, color=WHITE)


def make_scripts_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '07  校验与导出脚本', size=28, color=WHITE, bold=True)

    # Left: validation
    card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.9), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  校验脚本 (Validation Scripts)', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  执行时机：数据提交时自动执行', size=12, color=WHITE)
    add_para(tf, '  运行环境：Python 沙箱 (utils/script_runner.py)', size=12, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  脚本可访问的变量', size=13, color=ORANGE, bold=True)
    add_para(tf, '  record    — 当前提交的数据', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  errors    — 错误列表（阻止保存）', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  warnings  — 警告列表（允许保存）', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  relations — 关联操作队列', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  示例', size=13, color=ORANGE, bold=True)
    code = '''  # 校验开始日期不晚于结束日期
  start = record.get("startDate")
  end = record.get("endDate")
  if start and end and start > end:
      errors.append("开始日期不能晚于结束日期")

  # 校验名称长度
  name = record.get("caseName", "")
  if len(name) < 2:
      errors.append("名称至少2个字符")'''
    add_para(tf, code, size=10, color=ACCENT2, font_name='Consolas')

    # Right: export
    card2 = add_shape(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.9), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  导出脚本 (Export Scripts)', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  两种作用域', size=13, color=ORANGE, bold=True)
    add_para(tf, '  page  — 页面级导出（处理整个集合）', size=11, color=WHITE)
    add_para(tf, '  row   — 行级导出（处理单条记录）', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  脚本可访问的变量', size=13, color=ORANGE, bold=True)
    add_para(tf, '  records   — 数据列表', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  fields    — 字段配置', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  output    — 输出容器', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  支持输出格式', size=13, color=ORANGE, bold=True)
    add_para(tf, '  json | csv | xml | html | txt | xlsx', size=12, color=GREEN, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  绑定方式', size=13, color=ORANGE, bold=True)
    add_para(tf, '  page_configs.export_scripts', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  page_configs.row_export_scripts', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  同一脚本可绑定到多个页面（复用）', size=11, color=WHITE)


def make_etl_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '08  ETL 数据管道', size=28, color=WHITE, bold=True)

    card = add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(2.0), BG_CARD)
    tf = card.text_frame
    tf.word_wrap = True
    set_text(tf, '  ETL 管道：可视化步骤编排', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    flow = '  HTTP 数据抽取  →  Python 脚本转换  →  字段映射  →  写入目标集合  →  执行日志记录'
    add_para(tf, flow, size=14, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  每个 ETL 任务由多个 step 组成，存储在 etl_tasks.steps (JSONB)，执行结果记录到 etl_logs', size=12, color=GRAY)

    # Three step types
    types = [
        ('HTTP 抽取', 'httpExtract', [
            'url — 请求地址',
            'method — GET/POST',
            'headers — 请求头',
            'dataPath — 从响应 JSON 提取数据',
        ]),
        ('脚本转换', 'scriptTransform', [
            'Python 沙箱执行',
            '输入: records (上一步结果)',
            '输出: records (转换后)',
            '支持数据清洗、格式转换',
        ]),
        ('字段映射 + 写入', 'fieldMapping', [
            'targetCollection — 目标集合',
            'mappings — 源字段 → 目标字段',
            'mode — insert / upsert',
            '自动处理 ID 生成和去重',
        ]),
    ]
    for idx, (title, step_type, items) in enumerate(types):
        x = Inches(0.5 + idx * 4.2)
        card = add_shape(slide, x, Inches(3.5), Inches(3.9), Inches(3.6), BG_CARD)
        tf = card.text_frame
        tf.word_wrap = True
        set_text(tf, f'  {title}', size=15, color=ACCENT, bold=True)
        add_para(tf, f'  type: {step_type}', size=10, color=GRAY, font_name='Consolas')
        add_para(tf, '', size=4, color=WHITE)
        for item in items:
            add_para(tf, f'  • {item}', size=11, color=WHITE)


def make_openapi_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '09  Open API', size=28, color=WHITE, bold=True)

    # Left
    card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.9), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  API Key 认证机制', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  认证方式', size=13, color=ORANGE, bold=True)
    add_para(tf, '  Header: X-API-Key: <api_key>', size=12, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  密钥管理', size=13, color=ORANGE, bold=True)
    add_para(tf, '  • 密钥仅在创建时显示一次（SHA256 存储）', size=11, color=WHITE)
    add_para(tf, '  • 支持启用/禁用、记录最后使用时间', size=11, color=WHITE)
    add_para(tf, '  • 管理界面：系统配置 → 平台管理 → Open API', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  集合开放控制', size=13, color=ORANGE, bold=True)
    add_para(tf, '  page_configs.api_public = true', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  每个集合独立控制是否对 API 开放', size=11, color=WHITE)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  读写控制', size=13, color=ORANGE, bold=True)
    add_para(tf, '  page_configs.api_writable = true', size=11, color=ACCENT2, font_name='Consolas')
    add_para(tf, '  开启后支持 POST/PUT 写入操作', size=11, color=WHITE)

    # Right
    card2 = add_shape(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.9), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  API 端点', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    endpoints = [
        ('GET', '/api/open/<collection>', '查询列表 (支持 ?q= 查询)'),
        ('GET', '/api/open/<collection>/<id>', '获取单条记录'),
        ('POST', '/api/open/<collection>', '新增记录 (需 writable)'),
        ('PUT', '/api/open/<collection>/<id>', '更新记录 (需 writable)'),
    ]
    for method, path, desc in endpoints:
        color = GREEN if method == 'GET' else ORANGE
        add_para(tf, f'  {method}', size=12, color=color, bold=True, font_name='Consolas')
        add_para(tf, f'  {path}', size=10, color=ACCENT2, font_name='Consolas', space_before=Pt(0))
        add_para(tf, f'  {desc}', size=10, color=GRAY, space_before=Pt(0))
        add_para(tf, '', size=4, color=WHITE)

    add_para(tf, '  调用示例', size=13, color=ORANGE, bold=True)
    code = '''  curl -H "X-API-Key: ak_xxx" \\
    "http://host/api/open/inspection-case\\
    ?q={\\\"用例类型\\\":\\\"功能巡检\\\"}"'''
    add_para(tf, code, size=10, color=ACCENT2, font_name='Consolas')


def make_rbac_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '10  权限与审计', size=28, color=WHITE, bold=True)

    # Left: RBAC
    card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.9), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  RBAC 三级角色体系', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=6, color=WHITE)
    roles = [
        ('admin (管理员)', [
            '全部功能访问权限',
            '菜单/页面配置/用户管理',
            '系统配置、备份、导出脚本',
            '数据增删改查',
        ]),
        ('developer (开发者)', [
            '业务数据页面访问',
            '数据增删改查',
            '数据查询控制台',
            '不可访问系统配置',
        ]),
        ('guest (访客)', [
            '业务数据只读访问',
            '不可新增/修改/删除',
            '不可访问系统配置',
        ]),
    ]
    for role_name, perms in roles:
        color = GREEN if 'admin' in role_name else ACCENT if 'developer' in role_name else GRAY
        add_para(tf, f'  {role_name}', size=14, color=color, bold=True)
        for perm in perms:
            add_para(tf, f'     • {perm}', size=11, color=WHITE, space_before=Pt(1))
        add_para(tf, '', size=4, color=WHITE)

    # Right: audit
    card2 = add_shape(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.9), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  全量操作审计', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  记录内容 (operation_logs)', size=13, color=ORANGE, bold=True)
    fields = [
        'action      — create / update / delete',
        'target_type — 操作对象类型',
        'target_id   — 操作对象 ID',
        'target_name — 操作对象名称',
        'description — 变更描述（含字段级差异）',
        'operator    — 操作人信息',
        'batch_id    — 批量操作聚合标识',
        'created_at  — 操作时间',
    ]
    for f in fields:
        add_para(tf, f'  {f}', size=10, color=ACCENT2, font_name='Consolas')

    add_para(tf, '', size=6, color=WHITE)
    add_para(tf, '  权限实现方式', size=13, color=ORANGE, bold=True)
    add_para(tf, '  • menus.roles JSONB — 菜单可见性', size=11, color=WHITE)
    add_para(tf, '  • auth.py @login_required — 登录校验', size=11, color=WHITE)
    add_para(tf, '  • auth.py @admin_required — 管理员校验', size=11, color=WHITE)
    add_para(tf, '  • hasRoutePermission() — 前端路由守卫', size=11, color=WHITE)
    add_para(tf, '  • filterMenusByRole() — 前端菜单过滤', size=11, color=WHITE)


def make_backup_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '11  备份与恢复', size=28, color=WHITE, bold=True)

    card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.9), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  备份机制', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    features = [
        ('手动备份', '管理员在界面一键创建备份'),
        ('定时备份', 'backup_settings 配置间隔（daily/weekly/monthly）'),
        ('备份内容', '全量导出：menus + page_configs + dynamic_data + data_relations'),
        ('存储格式', 'JSON 文件，含元数据（表计数、记录数）'),
        ('文件管理', '支持下载备份文件到本地'),
        ('保留策略', 'retention_count 控制保留份数，自动清理旧备份'),
    ]
    for title, desc in features:
        add_para(tf, f'  {title}', size=13, color=ORANGE, bold=True)
        add_para(tf, f'     {desc}', size=11, color=WHITE, space_before=Pt(0))
        add_para(tf, '', size=2, color=WHITE)

    card2 = add_shape(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.9), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  恢复与对比', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  数据还原', size=13, color=ORANGE, bold=True)
    add_para(tf, '  选择备份版本 → 一键还原', size=11, color=WHITE)
    add_para(tf, '  清空当前数据后从备份文件恢复', size=11, color=WHITE)
    add_para(tf, '', size=6, color=WHITE)
    add_para(tf, '  数据对比 (Diff)', size=13, color=ORANGE, bold=True)
    add_para(tf, '  选择两个备份版本', size=11, color=WHITE)
    add_para(tf, '  逐集合、逐记录、逐字段对比', size=11, color=WHITE)
    add_para(tf, '  差异类型：新增 / 修改 / 删除', size=11, color=WHITE)
    add_para(tf, '  支持导出对比报告', size=11, color=WHITE)
    add_para(tf, '', size=6, color=WHITE)
    add_para(tf, '  跨环境迁移', size=13, color=ORANGE, bold=True)
    add_para(tf, '  下载备份 → 上传到目标环境 → 还原', size=11, color=WHITE)
    add_para(tf, '  适用于 开发 → 测试 → 生产 的数据同步', size=11, color=WHITE)


def make_deploy_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_text(tb.text_frame, '12  部署方案', size=28, color=WHITE, bold=True)

    card1 = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.2), Inches(5.9), BG_CARD)
    tf = card1.text_frame
    tf.word_wrap = True
    set_text(tf, '  本地部署', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    code = '''  # 一键启动
  ./start.sh

  # 启动流程:
  # 1. python init_db.py (数据库迁移)
  # 2. npm run build     (前端构建)
  # 3. python proxy.py   (反向代理+后端)

  # 服务架构:
  # 浏览器 → :8080 (proxy.py)
  #   ├── 静态文件 ← dist/
  #   └── /api/*   → Flask (:3001)
  #                    └── PostgreSQL'''
    add_para(tf, code, size=11, color=ACCENT2, font_name='Consolas')

    card2 = add_shape(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.9), BG_CARD)
    tf = card2.text_frame
    tf.word_wrap = True
    set_text(tf, '  公网访问 (frp 内网穿透)', size=16, color=ACCENT, bold=True)
    add_para(tf, '', size=4, color=WHITE)
    code2 = '''  # 服务器 (frps.toml)
  bindPort = 7000
  vhostHTTPPort = 80
  auth.token = "secret"

  # 本地 (frpc.toml)
  serverAddr = "your-server-ip"
  serverPort = 7000
  auth.token = "secret"

  [[proxies]]
  name = "check-manage"
  type = "http"
  localPort = 8080
  customDomains = ["check.example.com"]'''
    add_para(tf, code2, size=10, color=ACCENT2, font_name='Consolas')
    add_para(tf, '', size=4, color=WHITE)
    add_para(tf, '  两种模式', size=13, color=ORANGE, bold=True)
    add_para(tf, '  HTTP + 域名（推荐）', size=11, color=GREEN)
    add_para(tf, '  TCP + 端口（无需域名）', size=11, color=WHITE)


def make_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    tb = add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    set_text(tb.text_frame, '总结', size=32, color=WHITE, bold=True)

    highlights = [
        ('配置驱动', '新增业务实体零代码，INSERT 两行 SQL 即可获得完整 CRUD + UI + API'),
        ('JSONB 单表', '所有业务数据统一存储，无需 DDL 变更，灵活应对字段变化'),
        ('三种关联', 'relation（M:N 双向同步）/ reference（父子继承）/ quoteSelect（轻量引用）'),
        ('乐观锁', 'version 字段 + 双重校验，安全解决多人并发编辑冲突'),
        ('查询引擎', 'MongoDB 风格语法 → PostgreSQL JSONB，12 种操作符 + 连表查询'),
        ('脚本沙箱', 'Python 校验脚本 + 导出脚本 + ETL 转换，灵活可扩展'),
        ('Open API', 'API Key 认证，外部系统集成，按集合控制读写权限'),
        ('全链路审计', '操作日志 + 系统备份 + 数据对比，完整的数据生命周期管理'),
    ]
    y = Inches(1.5)
    for idx, (title, desc) in enumerate(highlights):
        col = idx % 2
        row = idx // 2
        x = Inches(0.5 + col * 6.4)
        cy = y + Inches(row * 1.35)
        card = add_shape(slide, x, cy, Inches(6.1), Inches(1.15), BG_CARD)
        tf = card.text_frame
        tf.word_wrap = True
        set_text(tf, f'  {title}', size=15, color=ACCENT, bold=True)
        add_para(tf, f'  {desc}', size=11, color=GRAY)

    # Bottom
    tb2 = add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4))
    set_text(tb2.text_frame, '详细文档：docs/ 目录下的语法手册、使用文档、优势分析报告、部署指南', size=12, color=GRAY, alignment=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)
    make_toc_slide(prs)
    make_arch_slide(prs)
    make_config_driven_slide(prs)
    make_controls_slide(prs)
    make_relation_slide(prs)
    make_conflict_slide(prs)
    make_query_slide(prs)
    make_scripts_slide(prs)
    make_etl_slide(prs)
    make_openapi_slide(prs)
    make_rbac_slide(prs)
    make_backup_slide(prs)
    make_deploy_slide(prs)
    make_summary_slide(prs)

    out = os.path.join(os.path.dirname(__file__), '..', 'docs', 'check-manage 系统技术介绍.pptx')
    prs.save(out)
    print(f'PPT saved: {out} ({len(prs.slides)} slides)')


if __name__ == '__main__':
    main()
