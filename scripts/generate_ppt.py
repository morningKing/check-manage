# -*- coding: utf-8 -*-
"""
动态数据管理平台 — 系统功能介绍 PPT 生成脚本
使用 python-pptx 生成专业现代风格演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math
import os

# ── 全局配色 ──────────────────────────────────────────────
PRIMARY = RGBColor(0x40, 0x9E, 0xFF)       # Element Plus 主色
PRIMARY_DARK = RGBColor(0x33, 0x7E, 0xCC)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # 深色背景
DARK_BG2 = RGBColor(0x16, 0x21, 0x3E)
ACCENT = RGBColor(0x67, 0xC2, 0x3A)        # 绿色强调
ACCENT2 = RGBColor(0xE6, 0xA2, 0x3C)       # 橙色强调
ACCENT3 = RGBColor(0xF5, 0x6C, 0x6C)       # 红色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xC0, 0xC4, 0xCC)
VERY_LIGHT = RGBColor(0xE4, 0xE7, 0xED)
CARD_BG = RGBColor(0x23, 0x2D, 0x4A)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # 圆角调整
    shape.adjustments[0] = 0.08
    return shape


def add_circle(slide, left, top, size, fill_color, alpha=None):
    """添加圆形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.rotation = 0
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_decorative_circles(slide):
    """添加装饰性圆形元素"""
    # 右上角大圆
    add_circle(slide, Inches(11.5), Inches(-0.8), Inches(2.5), RGBColor(0x40, 0x9E, 0xFF))
    # 左下角
    add_circle(slide, Inches(-0.5), Inches(5.8), Inches(2), RGBColor(0x67, 0xC2, 0x3A))
    # 右中小圆
    add_circle(slide, Inches(12.5), Inches(3.5), Inches(0.8), ACCENT2)


def add_top_bar(slide):
    """添加顶部装饰条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_page_number(slide, num, total):
    """添加页码"""
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num} / {total}", font_size=11, color=LIGHT_TEXT,
                 alignment=PP_ALIGN.RIGHT)


def add_feature_card(slide, left, top, width, height, icon_text, title, desc, icon_color=PRIMARY):
    """添加功能卡片"""
    card = add_rounded_rect(slide, left, top, width, height, CARD_BG,
                            border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

    # 图标圆
    icon_size = Inches(0.55)
    icon_left = left + Inches(0.3)
    icon_top = top + Inches(0.3)
    icon_circle = add_circle(slide, icon_left, icon_top, icon_size, icon_color)

    # 图标文字
    add_text_box(slide, icon_left, icon_top + Inches(0.08), icon_size, Inches(0.4),
                 icon_text, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 标题
    add_text_box(slide, left + Inches(0.3), top + Inches(1.05), width - Inches(0.6), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)

    # 描述
    add_text_box(slide, left + Inches(0.3), top + Inches(1.45), width - Inches(0.6), height - Inches(1.75),
                 desc, font_size=11, color=LIGHT_TEXT)

    return card


def build_slide_title(prs):
    """第1页：封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # 装饰圆
    add_circle(slide, Inches(10.0), Inches(-1.5), Inches(4.5), RGBColor(0x1E, 0x3A, 0x6E))
    add_circle(slide, Inches(11.5), Inches(0.0), Inches(3.0), RGBColor(0x24, 0x44, 0x7A))
    add_circle(slide, Inches(-1.5), Inches(5.0), Inches(3.5), RGBColor(0x1E, 0x3A, 0x6E))

    # 左侧色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), SLIDE_HEIGHT)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    # 主标题
    add_text_box(slide, Inches(1.2), Inches(1.8), Inches(8), Inches(1.2),
                 "动态数据管理平台", font_size=48, color=WHITE, bold=True)

    # 英文副标题
    add_text_box(slide, Inches(1.2), Inches(3.0), Inches(8), Inches(0.6),
                 "Dynamic Data Management Platform", font_size=22, color=PRIMARY)

    # 分割线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.2), Inches(3.7), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()

    # 描述
    add_text_box(slide, Inches(1.2), Inches(4.0), Inches(7), Inches(1.2),
                 "低代码驱动的企业级数据管理解决方案\n配置即开发，让数据管理更高效",
                 font_size=18, color=LIGHT_TEXT)

    # 技术标签
    tags = ["Vue 3", "TypeScript", "Element Plus", "Flask", "PostgreSQL"]
    x = Inches(1.2)
    for tag in tags:
        w = Inches(len(tag) * 0.15 + 0.4)
        add_rounded_rect(slide, x, Inches(5.4), w, Inches(0.38), RGBColor(0x2A, 0x3A, 0x5C),
                         border_color=PRIMARY, border_width=Pt(1))
        add_text_box(slide, x + Inches(0.05), Inches(5.42), w - Inches(0.1), Inches(0.35),
                     tag, font_size=11, color=PRIMARY, alignment=PP_ALIGN.CENTER)
        x += w + Inches(0.15)


def build_slide_overview(prs):
    """第2页：平台概述"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_top_bar(slide)
    add_page_number(slide, 2, 9)

    # 标题
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                 "平台概述", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8), Inches(0.5),
                 "一套配置驱动的低代码数据管理系统，让非开发人员也能快速构建数据应用",
                 font_size=14, color=LIGHT_TEXT)

    # 核心概念 — 3张卡片
    concepts = [
        ("配", "页面配置 PageConfig", "定义数据集合的字段架构\n字段名称、类型、验证规则\n以 JSONB 格式存储", PRIMARY),
        ("数", "动态数据 DynamicData", "所有业务数据统一存储\ncollection + data JSONB\n灵活适应任意数据结构", ACCENT),
        ("菜", "菜单管理 Menu", "将 URL 路径链接到页面配置\n动态生成侧边栏导航\n支持多级菜单树", ACCENT2),
    ]

    x = Inches(0.8)
    for icon, title, desc, color in concepts:
        card = add_rounded_rect(slide, x, Inches(1.8), Inches(3.7), Inches(2.4), CARD_BG,
                                border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

        # 图标
        ic = add_circle(slide, x + Inches(0.3), Inches(2.1), Inches(0.7), color)
        add_text_box(slide, x + Inches(0.3), Inches(2.22), Inches(0.7), Inches(0.5),
                     icon, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(1.2), Inches(2.15), Inches(2.3), Inches(0.4),
                     title, font_size=15, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), Inches(2.9), Inches(3.1), Inches(1.2),
                     desc, font_size=12, color=LIGHT_TEXT)

        x += Inches(3.95)

    # 工作流程
    add_text_box(slide, Inches(0.8), Inches(4.5), Inches(6), Inches(0.5),
                 "工作流程", font_size=20, color=WHITE, bold=True)

    steps = [
        ("01", "管理员配置", "定义字段架构\n和验证规则"),
        ("02", "系统自动生成", "CRUD 页面\n表单与表格"),
        ("03", "用户使用", "数据录入\n导入 / 导出"),
        ("04", "数据管理", "搜索、筛选\n关联、对比"),
    ]

    arrow_color = RGBColor(0x30, 0x3D, 0x5F)
    x = Inches(0.8)
    for i, (num, title, desc) in enumerate(steps):
        # 步骤卡片
        add_rounded_rect(slide, x, Inches(5.1), Inches(2.6), Inches(1.8), CARD_BG,
                         border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))
        # 编号
        add_text_box(slide, x + Inches(0.2), Inches(5.2), Inches(0.5), Inches(0.4),
                     num, font_size=22, color=PRIMARY, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(5.6), Inches(2.2), Inches(0.35),
                     title, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(5.95), Inches(2.2), Inches(0.8),
                     desc, font_size=11, color=LIGHT_TEXT)

        # 箭头（除最后一个）
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           x + Inches(2.7), Inches(5.85), Inches(0.35), Inches(0.25))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PRIMARY
            arrow.line.fill.background()

        x += Inches(3.05)


def build_slide_dynamic_page(prs):
    """第3页：动态页面配置"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_top_bar(slide)
    add_page_number(slide, 3, 9)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                 "动态页面配置", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                 "管理员通过可视化界面定义字段架构，系统自动生成完整的 CRUD 页面，无需编写代码",
                 font_size=14, color=LIGHT_TEXT)

    # 15种控件 — 分3行5列
    controls = [
        ("Aa", "文本 text", PRIMARY),
        ("¶", "多行 textarea", PRIMARY),
        ("#", "数字 number", PRIMARY),
        ("▾", "下拉 select", ACCENT),
        ("☑", "多选 multiSelect", ACCENT),
        ("◉", "单选 radio", ACCENT),
        ("☐", "复选 checkbox", ACCENT),
        ("📅", "日期 date", ACCENT2),
        ("⏰", "时间 datetime", ACCENT2),
        ("📎", "文件 file", ACCENT2),
        ("🖼", "图片 image", ACCENT2),
        ("🔗", "关联 relation", ACCENT3),
        ("↗", "引用 reference", ACCENT3),
        ("📌", "引选 quoteSelect", ACCENT3),
        ("⏱", "时间戳 auto", RGBColor(0x90, 0x9B, 0xF5)),
    ]

    x_start = Inches(0.8)
    y_start = Inches(1.7)
    card_w = Inches(2.3)
    card_h = Inches(0.8)
    gap_x = Inches(0.2)
    gap_y = Inches(0.15)

    for i, (icon, label, color) in enumerate(controls):
        col = i % 5
        row = i // 5
        x = x_start + col * (card_w + gap_x)
        y = y_start + row * (card_h + gap_y)

        add_rounded_rect(slide, x, y, card_w, card_h, CARD_BG,
                         border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

        # 颜色指示条
        indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x, y, Inches(0.06), card_h)
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = color
        indicator.line.fill.background()

        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.4), Inches(0.35),
                     icon, font_size=16, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.55), y + Inches(0.2), Inches(1.6), Inches(0.4),
                     label, font_size=12, color=WHITE, bold=True)

    # 右侧特性说明
    features_y = Inches(4.7)
    add_text_box(slide, Inches(0.8), features_y, Inches(6), Inches(0.5),
                 "页面配置核心能力", font_size=20, color=WHITE, bold=True)

    feature_items = [
        ("可视化字段编辑", "拖拽排序，实时预览表单效果"),
        ("字段验证规则", "必填、唯一、正则等多种验证"),
        ("选项数据源", "静态选项 / API接口 / 关联集合"),
        ("自动序号生成", "支持自定义前缀和编号规则"),
    ]

    x = Inches(0.8)
    for title, desc in feature_items:
        add_rounded_rect(slide, x, features_y + Inches(0.55), Inches(2.8), Inches(1.6), CARD_BG,
                         border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))
        # 小圆点
        dot = add_circle(slide, x + Inches(0.2), features_y + Inches(0.75), Inches(0.12), PRIMARY)
        add_text_box(slide, x + Inches(0.4), features_y + Inches(0.7), Inches(2.2), Inches(0.3),
                     title, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), features_y + Inches(1.1), Inches(2.4), Inches(0.8),
                     desc, font_size=11, color=LIGHT_TEXT)
        x += Inches(3.05)


def build_slide_relations(prs):
    """第4页：数据关系管理"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_top_bar(slide)
    add_page_number(slide, 4, 9)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                 "数据关系管理", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                 "三种灵活的数据关联模式，覆盖企业数据管理中的各类关系场景",
                 font_size=14, color=LIGHT_TEXT)

    # 三种关系类型卡片
    relations = [
        (
            "多对多关联", "Relation (M2M)",
            "⇄",
            PRIMARY,
            "双向关联，两端都可查看\n存储于 data_relations 表\n支持批量管理关联记录\n导入时自动按显示名匹配",
            "员工 ⇄ 项目\n标签 ⇄ 文章"
        ),
        (
            "父子引用", "Reference",
            "↓",
            ACCENT,
            "子记录存储父记录 ID\n支持字段继承 (inheritFields)\n父数据变更自动反映到子记录\n导入时按主键或显示名匹配",
            "部门 → 员工\n类别 → 商品"
        ),
        (
            "单向引用", "QuoteSelect",
            "→",
            ACCENT2,
            "单方向引用，存储 ID 数组\n可选择多条目标记录\n不影响被引用方的数据\n导入时按显示名自动解析",
            "订单 → 商品\n计划 → 任务"
        ),
    ]

    x = Inches(0.8)
    for title, subtitle, icon, color, desc, example in relations:
        card_w = Inches(3.7)
        card_h = Inches(5.0)
        add_rounded_rect(slide, x, Inches(1.7), card_w, card_h, CARD_BG,
                         border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

        # 顶部色带
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x, Inches(1.7), card_w, Inches(0.06))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()

        # 图标
        ic = add_circle(slide, x + Inches(1.35), Inches(2.0), Inches(1.0), color)
        add_text_box(slide, x + Inches(1.35), Inches(2.15), Inches(1.0), Inches(0.6),
                     icon, font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # 标题
        add_text_box(slide, x + Inches(0.3), Inches(3.15), Inches(3.1), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(3.1), Inches(0.3),
                     subtitle, font_size=12, color=color, alignment=PP_ALIGN.CENTER)

        # 分割线
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x + Inches(0.5), Inches(3.9), Inches(2.7), Inches(0.02))
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(0x30, 0x3D, 0x5F)
        sep.line.fill.background()

        # 描述
        add_text_box(slide, x + Inches(0.3), Inches(4.05), Inches(3.1), Inches(1.6),
                     desc, font_size=11, color=LIGHT_TEXT)

        # 示例
        add_text_box(slide, x + Inches(0.3), Inches(5.6), Inches(0.8), Inches(0.25),
                     "示例：", font_size=10, color=color, bold=True)
        add_text_box(slide, x + Inches(0.3), Inches(5.85), Inches(3.1), Inches(0.7),
                     example, font_size=10, color=VERY_LIGHT)

        x += Inches(3.95)


def build_slide_import_export(prs):
    """第5页：数据导入导出"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_top_bar(slide)
    add_page_number(slide, 5, 9)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                 "数据导入与导出", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                 "多格式文件支持，智能数据转换，让数据流转更加便捷",
                 font_size=14, color=LIGHT_TEXT)

    # 导入部分
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(3), Inches(0.5),
                 "📥  数据导入", font_size=20, color=PRIMARY, bold=True)

    import_card = add_rounded_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.6), CARD_BG,
                                   border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

    # 支持格式
    formats = [
        ("Excel", ".xlsx / .xls", PRIMARY),
        ("JSON", ".json（新增）", ACCENT),
    ]

    y = Inches(2.4)
    for fmt, ext, color in formats:
        add_rounded_rect(slide, Inches(1.1), y, Inches(2.2), Inches(0.6),
                         RGBColor(0x1A, 0x2A, 0x4A), border_color=color, border_width=Pt(1.5))
        add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(0.8), Inches(0.25),
                     fmt, font_size=13, color=color, bold=True)
        add_text_box(slide, Inches(1.2), y + Inches(0.3), Inches(2.0), Inches(0.25),
                     ext, font_size=10, color=LIGHT_TEXT)
        y += Inches(0.75)

    # 导入流程
    add_text_box(slide, Inches(1.1), Inches(4.0), Inches(5), Inches(0.4),
                 "智能导入流程", font_size=14, color=WHITE, bold=True)

    steps = [
        "① 文件解析 — Excel/JSON 统一转为记录数组",
        "② 字段映射 — 表头/key 自动匹配字段标签或标识",
        "③ 标签转值 — labelToValue 将显示文本转为存储值",
        "④ 关联解析 — 自动将关联名称解析为记录 ID",
        "⑤ 逐条保存 — 进度可视化，失败自动跳过",
    ]

    y = Inches(4.4)
    for step in steps:
        add_text_box(slide, Inches(1.3), y, Inches(4.8), Inches(0.32),
                     step, font_size=11, color=LIGHT_TEXT)
        y += Inches(0.35)

    # 导出部分
    add_text_box(slide, Inches(7.0), Inches(1.7), Inches(3), Inches(0.5),
                 "📤  数据导出", font_size=20, color=ACCENT2, bold=True)

    export_card = add_rounded_rect(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.6), CARD_BG,
                                   border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

    exports = [
        ("Excel 导出", "内置标准导出，值转标签\n自动处理关联和引用字段显示", PRIMARY),
        ("自定义脚本导出", "Python 脚本灵活定制\n支持页面级和行级导出", ACCENT2),
        ("模板下载", "一键下载 Excel 导入模板\n包含字段标题和选项说明", ACCENT),
    ]

    y = Inches(2.5)
    for title, desc, color in exports:
        # 左侧指示条
        indicator = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(7.3), y, Inches(0.06), Inches(0.9))
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = color
        indicator.line.fill.background()

        add_text_box(slide, Inches(7.6), y + Inches(0.05), Inches(4.5), Inches(0.3),
                     title, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, Inches(7.6), y + Inches(0.38), Inches(4.5), Inches(0.55),
                     desc, font_size=11, color=LIGHT_TEXT)
        y += Inches(1.15)

    # JSON 格式说明
    add_text_box(slide, Inches(7.3), Inches(5.2), Inches(5), Inches(0.3),
                 "JSON 导入格式示例", font_size=12, color=ACCENT, bold=True)

    json_example = '[\n  { "名称": "张三", "状态": "启用" },\n  { "标签": ["标签A", "标签B"] }\n]'
    code_bg = add_rounded_rect(slide, Inches(7.3), Inches(5.5), Inches(4.8), Inches(1.2),
                                RGBColor(0x0D, 0x13, 0x25))
    add_text_box(slide, Inches(7.5), Inches(5.55), Inches(4.5), Inches(1.1),
                 json_example, font_size=10, color=ACCENT, font_name="Consolas")


def build_slide_table(prs):
    """第6页：数据表格功能"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_top_bar(slide)
    add_page_number(slide, 6, 9)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                 "数据表格", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                 "功能丰富的数据表格组件，提供高效的数据浏览与操作体验",
                 font_size=14, color=LIGHT_TEXT)

    features = [
        ("🔍", "关键字搜索", "全字段模糊搜索\n实时过滤结果", PRIMARY),
        ("📊", "列筛选", "按列精确筛选\n数值/日期区间", ACCENT),
        ("📄", "分页浏览", "可配置每页条数\n快速翻页定位", ACCENT2),
        ("☑", "批量操作", "多选行操作\n批量删除", ACCENT3),
        ("👁", "查看详情", "全字段展示\n关联数据跳转", RGBColor(0x90, 0x9B, 0xF5)),
        ("✏", "行内编辑", "快捷编辑\n乐观锁保护", PRIMARY_DARK),
    ]

    x_start = Inches(0.8)
    y = Inches(1.7)
    card_w = Inches(3.7)
    card_h = Inches(1.5)

    for i, (icon, title, desc, color) in enumerate(features):
        col = i % 3
        row = i // 3
        x = x_start + col * (card_w + Inches(0.3))
        cy = y + row * (card_h + Inches(0.2))

        add_rounded_rect(slide, x, cy, card_w, card_h, CARD_BG,
                         border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

        ic = add_circle(slide, x + Inches(0.25), cy + Inches(0.25), Inches(0.6), color)
        add_text_box(slide, x + Inches(0.25), cy + Inches(0.32), Inches(0.6), Inches(0.4),
                     icon, font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(1.0), cy + Inches(0.25), Inches(2.5), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(1.0), cy + Inches(0.65), Inches(2.5), Inches(0.7),
                     desc, font_size=11, color=LIGHT_TEXT)

    # 数据对比
    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(6), Inches(0.5),
                 "数据对比功能", font_size=20, color=WHITE, bold=True)

    compare_card = add_rounded_rect(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.4), CARD_BG,
                                    border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

    compare_items = [
        ("备份快照对比", "将当前数据与历史备份逐条对比\n快速发现数据变更", PRIMARY),
        ("差异可视化", "新增 / 修改 / 删除分类标记\n变更字段高亮显示", ACCENT),
        ("管理员专属", "仅管理员可访问\n保障数据安全", ACCENT2),
    ]

    x = Inches(1.1)
    for title, desc, color in compare_items:
        dot = add_circle(slide, x, Inches(5.95), Inches(0.15), color)
        add_text_box(slide, x + Inches(0.25), Inches(5.9), Inches(3.2), Inches(0.3),
                     title, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.25), Inches(6.2), Inches(3.2), Inches(0.6),
                     desc, font_size=10, color=LIGHT_TEXT)
        x += Inches(3.9)


def build_slide_permission(prs):
    """第7页：权限管理与安全"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_top_bar(slide)
    add_page_number(slide, 7, 9)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                 "权限管理与审计", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                 "基于角色的访问控制 (RBAC) + 完整操作审计，保障数据安全合规",
                 font_size=14, color=LIGHT_TEXT)

    # 三种角色
    roles = [
        ("👑", "管理员 Admin", "完全访问权限",
         ["系统管理（菜单、页面配置、用户）", "数据读写和删除", "导入/导出全功能", "数据对比", "审计日志查看"],
         ACCENT2),
        ("🔧", "开发者 Developer", "配置范围读写",
         ["访问已授权的菜单页面", "数据读写操作", "导入/导出", "无法访问系统管理", "无法查看审计日志"],
         PRIMARY),
        ("👤", "访客 Guest", "只读访问",
         ["浏览已授权页面数据", "搜索和筛选功能", "无法新增/编辑/删除", "无法导入数据", "无法批量操作"],
         LIGHT_TEXT),
    ]

    x = Inches(0.8)
    for icon, title, subtitle, perms, color in roles:
        card_w = Inches(3.7)
        add_rounded_rect(slide, x, Inches(1.7), card_w, Inches(4.2), CARD_BG,
                         border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

        # 顶部色带
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x, Inches(1.7), card_w, Inches(0.06))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()

        add_text_box(slide, x + Inches(0.3), Inches(1.95), Inches(0.5), Inches(0.5),
                     icon, font_size=24, color=color)
        add_text_box(slide, x + Inches(0.85), Inches(1.95), Inches(2.5), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.85), Inches(2.3), Inches(2.5), Inches(0.3),
                     subtitle, font_size=11, color=color)

        # 权限列表
        y = Inches(2.8)
        for perm in perms:
            add_text_box(slide, x + Inches(0.3), y, Inches(3.1), Inches(0.3),
                         f"•  {perm}", font_size=11, color=LIGHT_TEXT)
            y += Inches(0.33)

        x += Inches(3.95)

    # 审计日志
    add_text_box(slide, Inches(0.8), Inches(6.1), Inches(6), Inches(0.4),
                 "📋  操作审计日志", font_size=18, color=WHITE, bold=True)

    audit_card = add_rounded_rect(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.8), CARD_BG,
                                  border_color=RGBColor(0x30, 0x3D, 0x5F), border_width=Pt(1))

    audit_items = ["记录所有数据变更操作", "操作人 / 时间 / 类型 / 详情",
                   "支持按时间范围查询", "批量操作统一标记 (Batch ID)"]
    x = Inches(1.1)
    for item in audit_items:
        add_text_box(slide, x, Inches(6.6), Inches(2.7), Inches(0.5),
                     f"✓  {item}", font_size=11, color=LIGHT_TEXT)
        x += Inches(2.85)


def build_slide_tech(prs):
    """第8页：技术架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_top_bar(slide)
    add_page_number(slide, 8, 9)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                 "技术架构", font_size=32, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                 "现代化全栈技术方案，前后端分离架构",
                 font_size=14, color=LIGHT_TEXT)

    # 前端
    add_text_box(slide, Inches(0.8), Inches(1.7), Inches(3), Inches(0.5),
                 "前端 Frontend", font_size=20, color=PRIMARY, bold=True)

    fe_card = add_rounded_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), CARD_BG,
                               border_color=PRIMARY, border_width=Pt(1.5))

    fe_items = [
        ("Vue 3", "Composition API + <script setup>", "响应式框架"),
        ("TypeScript", "完整类型系统，类型安全", "开发语言"),
        ("Pinia", "状态管理，模块化 Store", "状态管理"),
        ("Element Plus", "企业级 UI 组件库", "UI 框架"),
        ("Vite", "极速热更新，高效构建", "构建工具"),
        ("Vitest", "单元测试，jsdom 环境", "测试框架"),
    ]

    y = Inches(2.5)
    for name, desc, tag in fe_items:
        add_text_box(slide, Inches(1.1), y, Inches(1.5), Inches(0.3),
                     name, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, Inches(2.7), y, Inches(3.0), Inches(0.3),
                     desc, font_size=11, color=LIGHT_TEXT)

        tag_w = Inches(len(tag) * 0.14 + 0.3)
        add_rounded_rect(slide, Inches(5.5 - 0.3), y + Inches(0.02), tag_w, Inches(0.26),
                         RGBColor(0x1A, 0x2A, 0x4A), border_color=PRIMARY, border_width=Pt(0.75))
        add_text_box(slide, Inches(5.5 - 0.3), y + Inches(0.02), tag_w, Inches(0.26),
                     tag, font_size=8, color=PRIMARY, alignment=PP_ALIGN.CENTER)
        y += Inches(0.55)

    # 后端
    add_text_box(slide, Inches(7.0), Inches(1.7), Inches(3), Inches(0.5),
                 "后端 Backend", font_size=20, color=ACCENT, bold=True)

    be_card = add_rounded_rect(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5), CARD_BG,
                               border_color=ACCENT, border_width=Pt(1.5))

    be_items = [
        ("Python Flask", "轻量 Web 框架，13 个蓝图", "Web 框架"),
        ("PostgreSQL", "JSONB 动态数据存储", "数据库"),
        ("psycopg2", "连接池，高效数据库交互", "DB 驱动"),
        ("JWT Auth", "Token 认证 + 角色鉴权", "认证"),
        ("Pytest", "完整后端测试覆盖", "测试"),
        ("脚本引擎", "隔离命名空间执行验证/导出脚本", "扩展"),
    ]

    y = Inches(2.5)
    for name, desc, tag in be_items:
        add_text_box(slide, Inches(7.3), y, Inches(1.5), Inches(0.3),
                     name, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, Inches(8.9), y, Inches(3.0), Inches(0.3),
                     desc, font_size=11, color=LIGHT_TEXT)

        tag_w = Inches(len(tag) * 0.14 + 0.3)
        add_rounded_rect(slide, Inches(11.7 - 0.3), y + Inches(0.02), tag_w, Inches(0.26),
                         RGBColor(0x1A, 0x2A, 0x4A), border_color=ACCENT, border_width=Pt(0.75))
        add_text_box(slide, Inches(11.7 - 0.3), y + Inches(0.02), tag_w, Inches(0.26),
                     tag, font_size=8, color=ACCENT, alignment=PP_ALIGN.CENTER)
        y += Inches(0.55)

    # 数据库设计
    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(12), Inches(0.5),
                 "📦 核心数据表：dynamic_data (id TEXT PK, collection TEXT, data JSONB, _version INT)  —  所有业务数据统一存储",
                 font_size=12, color=LIGHT_TEXT)


def build_slide_end(prs):
    """第9页：结尾"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # 装饰
    add_circle(slide, Inches(10.0), Inches(-1.5), Inches(4.5), RGBColor(0x1E, 0x3A, 0x6E))
    add_circle(slide, Inches(11.5), Inches(0.0), Inches(3.0), RGBColor(0x24, 0x44, 0x7A))
    add_circle(slide, Inches(-1.5), Inches(5.0), Inches(3.5), RGBColor(0x1E, 0x3A, 0x6E))

    # 左侧色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), SLIDE_HEIGHT)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(10), Inches(1.0),
                 "谢谢", font_size=56, color=WHITE, bold=True)

    add_text_box(slide, Inches(1.2), Inches(3.3), Inches(10), Inches(0.5),
                 "Thank You", font_size=28, color=PRIMARY)

    # 分割线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.2), Inches(4.1), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()

    add_text_box(slide, Inches(1.2), Inches(4.4), Inches(8), Inches(0.8),
                 "动态数据管理平台  |  低代码  ·  高效率  ·  更安全",
                 font_size=16, color=LIGHT_TEXT)

    # 底部技术标签
    tags = ["Vue 3", "TypeScript", "Element Plus", "Flask", "PostgreSQL"]
    x = Inches(1.2)
    for tag in tags:
        w = Inches(len(tag) * 0.15 + 0.4)
        add_rounded_rect(slide, x, Inches(5.6), w, Inches(0.38), RGBColor(0x2A, 0x3A, 0x5C),
                         border_color=PRIMARY, border_width=Pt(1))
        add_text_box(slide, x + Inches(0.05), Inches(5.62), w - Inches(0.1), Inches(0.35),
                     tag, font_size=11, color=PRIMARY, alignment=PP_ALIGN.CENTER)
        x += w + Inches(0.15)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide_title(prs)       # 1. 封面
    build_slide_overview(prs)    # 2. 平台概述
    build_slide_dynamic_page(prs) # 3. 动态页面配置
    build_slide_relations(prs)   # 4. 数据关系管理
    build_slide_import_export(prs) # 5. 数据导入导出
    build_slide_table(prs)       # 6. 数据表格
    build_slide_permission(prs)  # 7. 权限管理
    build_slide_tech(prs)        # 8. 技术架构
    build_slide_end(prs)         # 9. 结尾

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "动态数据管理平台-系统功能介绍.pptx")
    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")


if __name__ == "__main__":
    main()
